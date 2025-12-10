"""
Ultimate Chat - A Modern AI Chat Interface
Multi-provider support: OpenAI, Anthropic, Google, xAI, Ollama
"""

import os
import json
import time
import uuid
import sqlite3
import traceback
import urllib.request
import urllib.error
import threading
from datetime import datetime
from pathlib import Path

from flask import Flask, request, Response, jsonify, send_from_directory, render_template_string
from dotenv import load_dotenv

load_dotenv()

# =============================================================================
# Task Manager for Cancellable Operations
# =============================================================================

class TaskManager:
    """Manages long-running tasks with cancellation support."""
    
    def __init__(self):
        self.tasks = {}  # task_id -> {status, progress, cancel_flag, result}
        self.lock = threading.Lock()
    
    def create_task(self, task_type):
        task_id = str(uuid.uuid4())[:8]
        with self.lock:
            self.tasks[task_id] = {
                'type': task_type,
                'status': 'starting',
                'progress': 0,
                'step': 0,
                'step_name': 'Initializing...',
                'cancel_requested': False,
                'result': None,
                'error': None,
                'created_at': time.time()
            }
        return task_id
    
    def update_task(self, task_id, **kwargs):
        with self.lock:
            if task_id in self.tasks:
                self.tasks[task_id].update(kwargs)
    
    def get_task(self, task_id):
        with self.lock:
            return self.tasks.get(task_id, {}).copy()
    
    def cancel_task(self, task_id):
        with self.lock:
            if task_id in self.tasks:
                self.tasks[task_id]['cancel_requested'] = True
                self.tasks[task_id]['status'] = 'cancelling'
                return True
        return False
    
    def is_cancelled(self, task_id):
        with self.lock:
            return self.tasks.get(task_id, {}).get('cancel_requested', False)
    
    def cleanup_old_tasks(self, max_age=3600):
        """Remove tasks older than max_age seconds."""
        now = time.time()
        with self.lock:
            to_remove = [tid for tid, t in self.tasks.items() if now - t['created_at'] > max_age]
            for tid in to_remove:
                del self.tasks[tid]

task_manager = TaskManager()

# =============================================================================
# Configuration
# =============================================================================

APP_PORT = int(os.getenv("PORT", os.getenv("APP_PORT", 5001)))  # Render uses PORT
MAX_UPLOAD_SIZE = int(os.getenv("MAX_UPLOAD_SIZE", 100)) * 1024 * 1024
ENABLE_IMAGE_GEN = os.getenv("ENABLE_IMAGE_GENERATION", "false").lower() == "true"
ENABLE_VIDEO_GEN = os.getenv("ENABLE_VIDEO_GENERATION", "false").lower() == "true"
IS_PRODUCTION = os.getenv("RENDER", "false").lower() == "true" or os.getenv("PRODUCTION", "false").lower() == "true"

# Paths - Use /data for Render persistent disk, otherwise local workspace
BASE_DIR = Path(__file__).parent
if IS_PRODUCTION and Path("/data").exists():
    DATA_DIR = Path("/data")
else:
    DATA_DIR = BASE_DIR

WORKSPACE_DIR = DATA_DIR / "workspace"
UPLOAD_DIR = WORKSPACE_DIR / "uploads"
GALLERY_DIR = WORKSPACE_DIR / "gallery"
VIDEO_DIR = WORKSPACE_DIR / "videos"
DOCUMENT_DIR = WORKSPACE_DIR / "documents"
DB_PATH = DATA_DIR / "chat.db"

for dir_path in [UPLOAD_DIR, GALLERY_DIR, VIDEO_DIR, DOCUMENT_DIR]:
    dir_path.mkdir(parents=True, exist_ok=True)

SYSTEM_PROMPT = """You are Ultimate Chat, a helpful, creative, and intelligent AI assistant. 
You can help with coding, writing, analysis, and creative tasks.
Be concise but thorough in your responses. Use markdown formatting when appropriate."""

# Provider configurations - Updated with latest models (June 2025)
PROVIDERS = {
    "openai": {
        "name": "OpenAI",
        "base_url": "https://api.openai.com/v1/chat/completions",
        "models": [
            "gpt-4.1", "gpt-4.1-mini", "gpt-4.1-nano",  # Latest GPT-4.1 series
            "gpt-4o", "gpt-4o-mini",  # GPT-4o series
            "o3", "o3-mini",  # Latest reasoning models
            "o1", "o1-mini", "o1-pro",  # O1 reasoning series
            "gpt-4-turbo", "gpt-4", "gpt-3.5-turbo"  # Legacy models
        ],
        "icon": "üü¢"
    },
    "anthropic": {
        "name": "Anthropic",
        "base_url": "https://api.anthropic.com/v1/messages",
        "models": [
            "claude-opus-4-20250514", "claude-sonnet-4-20250514",  # Claude 4 series
            "claude-3-7-sonnet-20250219",  # Claude 3.7
            "claude-3-5-sonnet-20241022", "claude-3-5-haiku-20241022",  # Claude 3.5 series
            "claude-3-opus-20240229", "claude-3-sonnet-20240229", "claude-3-haiku-20240307"  # Claude 3 series
        ],
        "icon": "üü§"
    },
    "google": {
        "name": "Google",
        "base_url": "https://generativelanguage.googleapis.com/v1beta/models/{model}:streamGenerateContent",
        "models": [
            "gemini-2.5-pro-preview-06-05",  # Latest Gemini 2.5
            "gemini-2.5-flash-preview-05-20",  # Gemini 2.5 Flash
            "gemini-2.0-flash", "gemini-2.0-flash-lite",  # Gemini 2.0 series
            "gemini-1.5-pro", "gemini-1.5-flash", "gemini-1.5-flash-8b"  # Gemini 1.5 series
        ],
        "icon": "üîµ"
    },
    "xai": {
        "name": "xAI",
        "base_url": "https://api.x.ai/v1/chat/completions",
        "models": [
            "grok-3", "grok-3-fast",  # Latest Grok 3
            "grok-2", "grok-2-vision",  # Grok 2 series
            "grok-beta"  # Beta model
        ],
        "icon": "‚ö´"
    },
    "ollama": {
        "name": "Ollama",
        "base_url": "http://localhost:11434/api/chat",
        "models": [
            "llama3.3:latest", "llama3.2:latest", "llama3.1:latest",  # Llama 3 series
            "qwen2.5:latest", "qwen2.5-coder:latest",  # Qwen series
            "deepseek-r1:latest", "deepseek-coder-v2:latest",  # DeepSeek
            "mistral:latest", "mixtral:latest",  # Mistral
            "codellama:latest", "phi4:latest", "gemma2:latest"  # Other popular models
        ],
        "icon": "ü¶ô"
    }
}

# =============================================================================
# Flask App Setup
# =============================================================================

app = Flask(__name__)
app.config['MAX_CONTENT_LENGTH'] = MAX_UPLOAD_SIZE

# =============================================================================
# Database Setup
# =============================================================================

def get_db():
    conn = sqlite3.connect(str(DB_PATH))
    conn.row_factory = sqlite3.Row
    return conn

def init_db():
    conn = get_db()
    conn.executescript('''
        CREATE TABLE IF NOT EXISTS conversations (
            id TEXT PRIMARY KEY,
            title TEXT,
            created_at TEXT,
            updated_at TEXT,
            message_count INTEGER DEFAULT 0
        );
        
        CREATE TABLE IF NOT EXISTS messages (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            conversation_id TEXT,
            role TEXT,
            content TEXT,
            attachments TEXT,
            created_at TEXT,
            FOREIGN KEY (conversation_id) REFERENCES conversations(id) ON DELETE CASCADE
        );
        
        CREATE TABLE IF NOT EXISTS api_keys (
            id TEXT PRIMARY KEY,
            provider TEXT UNIQUE,
            api_key TEXT,
            base_url TEXT,
            created_at TEXT,
            updated_at TEXT
        );
        
        CREATE TABLE IF NOT EXISTS custom_models (
            id TEXT PRIMARY KEY,
            provider TEXT,
            model_id TEXT,
            display_name TEXT,
            created_at TEXT
        );
        
        CREATE TABLE IF NOT EXISTS settings (
            key TEXT PRIMARY KEY,
            value TEXT
        );
        
        CREATE TABLE IF NOT EXISTS generated_files (
            id TEXT PRIMARY KEY,
            conversation_id TEXT,
            file_type TEXT,
            file_name TEXT,
            file_path TEXT,
            file_size INTEGER,
            mime_type TEXT,
            prompt TEXT,
            metadata TEXT,
            created_at TEXT,
            FOREIGN KEY (conversation_id) REFERENCES conversations(id) ON DELETE SET NULL
        );
        
        CREATE INDEX IF NOT EXISTS idx_files_type ON generated_files(file_type);
        CREATE INDEX IF NOT EXISTS idx_files_created ON generated_files(created_at);
    ''')
    conn.commit()
    conn.close()

# =============================================================================
# Generated Files Manager
# =============================================================================

def save_generated_file(file_type, file_name, file_path, prompt="", conversation_id=None, metadata=None):
    """Save a generated file record to the database."""
    conn = get_db()
    file_id = str(uuid.uuid4())[:12]
    
    # Get file size
    full_path = BASE_DIR / file_path.lstrip('/')
    file_size = full_path.stat().st_size if full_path.exists() else 0
    
    # Determine mime type
    mime_types = {
        'image': 'image/png',
        'audio': 'audio/mpeg',
        'video': 'video/mp4',
        'pdf': 'application/pdf',
        'docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
        'pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
        'xlsx': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
        'qr': 'image/png',
        'chart': 'image/png'
    }
    
    conn.execute('''INSERT INTO generated_files 
                    (id, conversation_id, file_type, file_name, file_path, file_size, mime_type, prompt, metadata, created_at)
                    VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?)''',
                 (file_id, conversation_id, file_type, file_name, file_path, file_size, 
                  mime_types.get(file_type, 'application/octet-stream'), prompt, 
                  json.dumps(metadata) if metadata else '{}', datetime.now().isoformat()))
    conn.commit()
    conn.close()
    return file_id

def get_generated_files(file_type=None, limit=50):
    """Get list of generated files."""
    conn = get_db()
    if file_type:
        rows = conn.execute('SELECT * FROM generated_files WHERE file_type = ? ORDER BY created_at DESC LIMIT ?', 
                           (file_type, limit)).fetchall()
    else:
        rows = conn.execute('SELECT * FROM generated_files ORDER BY created_at DESC LIMIT ?', (limit,)).fetchall()
    conn.close()
    return [dict(row) for row in rows]

def delete_generated_file(file_id):
    """Delete a generated file and its record."""
    conn = get_db()
    row = conn.execute('SELECT file_path FROM generated_files WHERE id = ?', (file_id,)).fetchone()
    if row:
        # Delete physical file
        full_path = BASE_DIR / row['file_path'].lstrip('/')
        if full_path.exists():
            full_path.unlink()
        # Delete record
        conn.execute('DELETE FROM generated_files WHERE id = ?', (file_id,))
        conn.commit()
    conn.close()

init_db()

# =============================================================================
# Settings Manager
# =============================================================================

def get_setting(key, default=None):
    conn = get_db()
    row = conn.execute('SELECT value FROM settings WHERE key = ?', (key,)).fetchone()
    conn.close()
    return row['value'] if row else default

def generate_smart_title(message):
    """Generate a smart, concise title from a message."""
    import re
    
    # Clean the message
    text = message.strip()
    
    # Remove common prefixes
    prefixes_to_remove = [
        r'^(hey|hi|hello|please|can you|could you|would you|i want to|i need to|help me|tell me about|what is|what are|how to|how do i|explain)\s*',
    ]
    for pattern in prefixes_to_remove:
        text = re.sub(pattern, '', text, flags=re.IGNORECASE)
    
    # If it's a question, extract the topic
    if '?' in text:
        text = text.split('?')[0]
    
    # If it's a command for generation, make it descriptive
    gen_patterns = [
        (r'generate\s+(image|video|audio|pdf|chart|qr)', r'\1 generation'),
        (r'create\s+(image|video|audio|pdf|chart|qr)', r'\1 creation'),
        (r'make\s+(a|an)?\s*(image|video|audio|pdf|chart|qr)', r'\2 generation'),
    ]
    for pattern, replacement in gen_patterns:
        if re.search(pattern, text, re.IGNORECASE):
            # Get what comes after
            match = re.search(pattern + r'\s*(.*)', text, re.IGNORECASE)
            if match:
                topic = match.group(2) if len(match.groups()) > 1 else match.group(1)
                subject = match.group(3) if len(match.groups()) > 2 else ''
                if subject:
                    text = f"{topic.title()}: {subject[:30]}"
                else:
                    text = f"{topic.title()} Generation"
                break
    
    # Capitalize first letter of each word for title case
    words = text.split()
    
    # Limit to key words (first 6 words or 40 chars)
    title_words = []
    char_count = 0
    for word in words[:8]:
        if char_count + len(word) > 40:
            break
        title_words.append(word)
        char_count += len(word) + 1
    
    title = ' '.join(title_words)
    
    # Title case, but keep small words lowercase
    small_words = {'a', 'an', 'the', 'and', 'or', 'but', 'in', 'on', 'at', 'to', 'for', 'of', 'with', 'is', 'are'}
    words = title.split()
    if words:
        words[0] = words[0].capitalize()
        for i in range(1, len(words)):
            if words[i].lower() not in small_words:
                words[i] = words[i].capitalize()
        title = ' '.join(words)
    
    # Fallback if title is too short
    if len(title) < 3:
        title = message[:40] + ('...' if len(message) > 40 else '')
    
    return title

def set_setting(key, value):
    conn = get_db()
    conn.execute('INSERT OR REPLACE INTO settings (key, value) VALUES (?, ?)', (key, value))
    conn.commit()
    conn.close()

def get_api_key(provider):
    conn = get_db()
    row = conn.execute('SELECT api_key, base_url FROM api_keys WHERE provider = ?', (provider,)).fetchone()
    conn.close()
    return dict(row) if row else None

def save_api_key(provider, api_key, base_url=None):
    conn = get_db()
    now = datetime.now().isoformat()
    if base_url is None:
        base_url = PROVIDERS.get(provider, {}).get('base_url', '')
    conn.execute('''INSERT OR REPLACE INTO api_keys (id, provider, api_key, base_url, created_at, updated_at) 
                    VALUES (?, ?, ?, ?, ?, ?)''',
                 (str(uuid.uuid4()), provider, api_key, base_url, now, now))
    conn.commit()
    conn.close()

def get_all_api_keys():
    conn = get_db()
    rows = conn.execute('SELECT provider, api_key, base_url FROM api_keys').fetchall()
    conn.close()
    return {row['provider']: {'api_key': row['api_key'], 'base_url': row['base_url']} for row in rows}

def get_custom_models():
    conn = get_db()
    rows = conn.execute('SELECT * FROM custom_models ORDER BY created_at DESC').fetchall()
    conn.close()
    return [dict(row) for row in rows]

def add_custom_model(provider, model_id, display_name=None):
    conn = get_db()
    conn.execute('INSERT INTO custom_models (id, provider, model_id, display_name, created_at) VALUES (?, ?, ?, ?, ?)',
                 (str(uuid.uuid4()), provider, model_id, display_name or model_id, datetime.now().isoformat()))
    conn.commit()
    conn.close()

def delete_custom_model(model_id):
    conn = get_db()
    conn.execute('DELETE FROM custom_models WHERE id = ?', (model_id,))
    conn.commit()
    conn.close()

# =============================================================================
# Multi-Provider Streaming Client
# =============================================================================

def stream_openai(messages, model, api_key, base_url=None):
    """Stream from OpenAI-compatible API."""
    url = base_url or "https://api.openai.com/v1/chat/completions"
    payload = {"model": model, "messages": messages, "stream": True}
    headers = {"Content-Type": "application/json", "Authorization": f"Bearer {api_key}"}
    
    req = urllib.request.Request(url, json.dumps(payload).encode('utf-8'), headers, method='POST')
    
    try:
        response = urllib.request.urlopen(req, timeout=120)
    except urllib.error.HTTPError as e:
        if e.code == 429:
            yield "\n\n‚ö†Ô∏è **Rate Limited by OpenAI.** Wait a moment and try again, or switch to another provider."
            return
        elif e.code == 401:
            yield "\n\n‚ùå **Invalid OpenAI API Key.** Please check your API key in Settings ‚öôÔ∏è"
            return
        else:
            yield f"\n\n‚ùå **OpenAI API Error ({e.code}):** {e.reason}"
            return
    except Exception as e:
        yield f"\n\n‚ùå **Connection Error:** {str(e)}"
        return
    
    with response:
        buffer = ""
        for chunk in iter(lambda: response.read(1024), b''):
            buffer += chunk.decode('utf-8')
            while '\n' in buffer:
                line, buffer = buffer.split('\n', 1)
                if line.startswith('data: ') and line != 'data: [DONE]':
                    try:
                        data = json.loads(line[6:])
                        if data.get('choices') and data['choices'][0].get('delta', {}).get('content'):
                            yield data['choices'][0]['delta']['content']
                    except: pass

def stream_anthropic(messages, model, api_key):
    """Stream from Anthropic API."""
    url = "https://api.anthropic.com/v1/messages"
    system = next((m['content'] for m in messages if m['role'] == 'system'), "")
    msgs = [m for m in messages if m['role'] != 'system']
    
    payload = {"model": model, "max_tokens": 8192, "stream": True, "messages": msgs}
    if system: payload["system"] = system
    
    headers = {"Content-Type": "application/json", "x-api-key": api_key, "anthropic-version": "2023-06-01"}
    req = urllib.request.Request(url, json.dumps(payload).encode('utf-8'), headers, method='POST')
    
    try:
        response = urllib.request.urlopen(req, timeout=120)
    except urllib.error.HTTPError as e:
        if e.code == 429:
            yield "\n\n‚ö†Ô∏è **Rate Limited by Anthropic.** Wait a moment and try again, or switch to another provider."
            return
        elif e.code == 401:
            yield "\n\n‚ùå **Invalid Anthropic API Key.** Please check your API key in Settings ‚öôÔ∏è"
            return
        else:
            yield f"\n\n‚ùå **Anthropic API Error ({e.code}):** {e.reason}"
            return
    except Exception as e:
        yield f"\n\n‚ùå **Connection Error:** {str(e)}"
        return
    
    with response:
        buffer = ""
        for chunk in iter(lambda: response.read(1024), b''):
            buffer += chunk.decode('utf-8')
            while '\n' in buffer:
                line, buffer = buffer.split('\n', 1)
                if line.startswith('data: '):
                    try:
                        data = json.loads(line[6:])
                        if data.get('type') == 'content_block_delta':
                            yield data.get('delta', {}).get('text', '')
                    except: pass

def stream_google(messages, model, api_key):
    """Stream from Google Gemini API."""
    url = f"https://generativelanguage.googleapis.com/v1beta/models/{model}:streamGenerateContent?key={api_key}&alt=sse"
    
    contents = []
    system_instruction = None
    for m in messages:
        if m['role'] == 'system':
            system_instruction = {"parts": [{"text": m['content']}]}
        else:
            role = "user" if m['role'] == 'user' else "model"
            contents.append({"role": role, "parts": [{"text": m['content']}]})
    
    payload = {"contents": contents}
    if system_instruction: payload["system_instruction"] = system_instruction
    
    headers = {"Content-Type": "application/json"}
    req = urllib.request.Request(url, json.dumps(payload).encode('utf-8'), headers, method='POST')
    
    try:
        response = urllib.request.urlopen(req, timeout=120)
    except urllib.error.HTTPError as e:
        if e.code == 429:
            yield "\n\n‚ö†Ô∏è **Rate Limited by Google.** Try switching to another provider (OpenAI, Claude, or Ollama) in the model selector."
            return
        elif e.code == 401 or e.code == 403:
            yield "\n\n‚ùå **Invalid Google API Key.** Please check your API key in Settings ‚öôÔ∏è"
            return
        else:
            yield f"\n\n‚ùå **Google API Error ({e.code}):** {e.reason}"
            return
    except Exception as e:
        yield f"\n\n‚ùå **Connection Error:** {str(e)}"
        return
    
    with response:
        buffer = ""
        for chunk in iter(lambda: response.read(1024), b''):
            buffer += chunk.decode('utf-8')
            while '\n' in buffer:
                line, buffer = buffer.split('\n', 1)
                if line.startswith('data: '):
                    try:
                        data = json.loads(line[6:])
                        if data.get('candidates'):
                            for part in data['candidates'][0].get('content', {}).get('parts', []):
                                if part.get('text'): yield part['text']
                    except: pass

def stream_xai(messages, model, api_key):
    """Stream from xAI Grok API."""
    yield from stream_openai(messages, model, api_key, "https://api.x.ai/v1/chat/completions")

def check_ollama_running(base_url=None):
    """Check if Ollama is running and accessible."""
    url = (base_url or "http://localhost:11434") + "/api/tags"
    try:
        req = urllib.request.Request(url, method='GET')
        with urllib.request.urlopen(req, timeout=5) as response:
            data = json.loads(response.read().decode('utf-8'))
            return True, [m['name'] for m in data.get('models', [])]
    except Exception as e:
        return False, str(e)

def stream_ollama(messages, model, base_url=None):
    """Stream from Ollama API with better error handling."""
    base = base_url or "http://localhost:11434"
    
    # First check if Ollama is running
    is_running, result = check_ollama_running(base)
    if not is_running:
        yield f"\n\n‚ö†Ô∏è **Ollama is not running or not accessible.**\n\nPlease:\n1. Install Ollama from https://ollama.ai\n2. Start Ollama\n3. Pull a model: `ollama pull {model}`\n\nOr switch to a cloud provider (OpenAI, Claude, etc.) in Settings ‚öôÔ∏è"
        return
    
    # Check if model exists
    available_models = result
    model_base = model.split(':')[0]
    if not any(model_base in m for m in available_models):
        yield f"\n\n‚ö†Ô∏è **Model '{model}' not found in Ollama.**\n\nAvailable models: {', '.join(available_models[:5]) if available_models else 'None'}\n\nPull it with: `ollama pull {model}`"
        return
    
    url = f"{base}/api/chat"
    payload = {"model": model, "messages": messages, "stream": True}
    headers = {"Content-Type": "application/json"}
    
    try:
        req = urllib.request.Request(url, json.dumps(payload).encode('utf-8'), headers, method='POST')
        with urllib.request.urlopen(req, timeout=180) as response:  # Increased timeout
            buffer = ""
            for chunk in iter(lambda: response.read(1024), b''):
                buffer += chunk.decode('utf-8')
                while '\n' in buffer:
                    line, buffer = buffer.split('\n', 1)
                    if line.strip():
                        try:
                            data = json.loads(line)
                            if data.get('message', {}).get('content'):
                                yield data['message']['content']
                        except: pass
    except urllib.error.URLError as e:
        yield f"\n\n‚ùå **Connection error:** {str(e.reason)}\n\nMake sure Ollama is running."
    except Exception as e:
        yield f"\n\n‚ùå **Error:** {str(e)}"

def stream_response(provider, model, messages):
    """Route to appropriate provider streaming function."""
    creds = get_api_key(provider)
    api_key = creds['api_key'] if creds else None
    base_url = creds.get('base_url') if creds else None
    
    try:
        if provider == "openai":
            if not api_key: raise ValueError("OpenAI API key not configured")
            yield from stream_openai(messages, model, api_key, base_url if base_url != PROVIDERS['openai']['base_url'] else None)
        elif provider == "anthropic":
            if not api_key: raise ValueError("Anthropic API key not configured")
            yield from stream_anthropic(messages, model, api_key)
        elif provider == "google":
            if not api_key: raise ValueError("Google API key not configured")
            yield from stream_google(messages, model, api_key)
        elif provider == "xai":
            if not api_key: raise ValueError("xAI API key not configured")
            yield from stream_xai(messages, model, api_key)
        elif provider == "ollama":
            yield from stream_ollama(messages, model, base_url)
        else:
            raise ValueError(f"Unknown provider: {provider}")
    except urllib.error.HTTPError as e:
        error_body = e.read().decode('utf-8') if e.fp else str(e)
        print(f"[{provider}] HTTP {e.code}: {error_body}")
        yield f"\n\n*[Error from {provider}: {e.code} - {e.reason}]*"
    except Exception as e:
        print(f"[{provider}] Error: {e}")
        traceback.print_exc()
        yield f"\n\n*[Error: {str(e)}]*"

# =============================================================================
# Image Generation (Stable Diffusion)
# =============================================================================

sd_pipeline = None
svd_pipeline = None

# GPU Settings from env
SD_MODEL = os.getenv("SD_MODEL", "stabilityai/stable-diffusion-xl-base-1.0")
SD_STEPS = int(os.getenv("SD_STEPS", 30))
SD_GUIDANCE = float(os.getenv("SD_GUIDANCE_SCALE", 7.5))
VIDEO_MODEL = os.getenv("VIDEO_MODEL", "stabilityai/stable-video-diffusion-img2vid-xt")
VIDEO_FRAMES = int(os.getenv("VIDEO_FRAMES", 25))
VIDEO_FPS = int(os.getenv("VIDEO_FPS", 7))

def load_sd_pipeline():
    """Load Stable Diffusion XL pipeline optimized for RTX 4090."""
    global sd_pipeline
    if sd_pipeline is not None:
        return sd_pipeline
    
    if not ENABLE_IMAGE_GEN:
        return None
    
    try:
        import torch
        
        # Check CUDA availability
        if not torch.cuda.is_available():
            print("[SD] CUDA not available, image generation disabled")
            return None
        
        gpu_name = torch.cuda.get_device_name(0)
        vram = torch.cuda.get_device_properties(0).total_memory / 1024**3
        print(f"[SD] Detected GPU: {gpu_name} ({vram:.1f} GB VRAM)")
        
        # Use SDXL for better quality on 4090
        if "stabilityai/stable-diffusion-xl" in SD_MODEL:
            from diffusers import StableDiffusionXLPipeline, AutoencoderKL
            
            print(f"[SD] Loading SDXL model: {SD_MODEL}")
            
            # Load VAE for better quality
            vae = AutoencoderKL.from_pretrained(
                "madebyollin/sdxl-vae-fp16-fix",
                torch_dtype=torch.float16
            )
            
            sd_pipeline = StableDiffusionXLPipeline.from_pretrained(
                SD_MODEL,
                vae=vae,
                torch_dtype=torch.float16,
                variant="fp16",
                use_safetensors=True
            ).to("cuda")
            
        else:
            from diffusers import StableDiffusionPipeline
            
            print(f"[SD] Loading SD model: {SD_MODEL}")
            sd_pipeline = StableDiffusionPipeline.from_pretrained(
                SD_MODEL,
                torch_dtype=torch.float16,
                safety_checker=None
            ).to("cuda")
        
        # Enable memory optimizations for 4090
        try:
            sd_pipeline.enable_xformers_memory_efficient_attention()
            print("[SD] xformers memory efficient attention enabled")
        except:
            print("[SD] xformers not available, using default attention")
        
        # Enable VAE slicing for large images
        sd_pipeline.enable_vae_slicing()
        
        print(f"[SD] Pipeline loaded successfully! Ready to generate images.")
        return sd_pipeline
        
    except Exception as e:
        print(f"[SD] Failed to load pipeline: {e}")
        traceback.print_exc()
        return None

def load_video_pipeline():
    """Load Stable Video Diffusion pipeline for RTX 4090."""
    global svd_pipeline
    if svd_pipeline is not None:
        return svd_pipeline
    
    if not ENABLE_VIDEO_GEN:
        return None
    
    try:
        import torch
        from diffusers import StableVideoDiffusionPipeline
        from diffusers.utils import export_to_video
        
        if not torch.cuda.is_available():
            print("[SVD] CUDA not available, video generation disabled")
            return None
        
        print(f"[SVD] Loading video model: {VIDEO_MODEL}")
        
        svd_pipeline = StableVideoDiffusionPipeline.from_pretrained(
            VIDEO_MODEL,
            torch_dtype=torch.float16,
            variant="fp16"
        ).to("cuda")
        
        # Enable memory optimizations
        try:
            svd_pipeline.enable_xformers_memory_efficient_attention()
        except:
            pass
        
        print("[SVD] Video pipeline loaded successfully!")
        return svd_pipeline
        
    except Exception as e:
        print(f"[SVD] Failed to load pipeline: {e}")
        traceback.print_exc()
        return None

def generate_image_with_progress(prompt, task_id, num_steps=None):
    """Generate image with SDXL and progress tracking."""
    import torch
    
    num_steps = num_steps or SD_STEPS
    
    def progress_callback(pipe, step, timestep, callback_kwargs):
        if task_manager.is_cancelled(task_id):
            raise InterruptedError("Generation cancelled by user")
        progress = int(10 + (step / num_steps) * 80)  # 10-90% for generation
        task_manager.update_task(task_id, 
            step=2, 
            step_name=f'Generating image... Step {step}/{num_steps}',
            progress=progress,
            status='running'
        )
        return callback_kwargs
    
    task_manager.update_task(task_id, step=1, step_name='Loading Stable Diffusion XL...', progress=5)
    
    pipeline = load_sd_pipeline()
    if pipeline is None:
        task_manager.update_task(task_id, status='error', error='Image generation not available. Check CUDA/GPU setup.')
        return None, "Image generation not available"
    
    if task_manager.is_cancelled(task_id):
        task_manager.update_task(task_id, status='cancelled')
        return None, "Cancelled"
    
    try:
        task_manager.update_task(task_id, step=2, step_name='Starting generation...', progress=10)
        
        # Generate with SDXL settings
        with torch.inference_mode():
            result = pipeline(
                prompt=prompt,
                num_inference_steps=num_steps,
                guidance_scale=SD_GUIDANCE,
                height=1024,
                width=1024,
                callback_on_step_end=progress_callback,
            )
        
        image = result.images[0]
        
        if task_manager.is_cancelled(task_id):
            task_manager.update_task(task_id, status='cancelled')
            return None, "Cancelled"
        
        task_manager.update_task(task_id, step=3, step_name='Saving to gallery...', progress=95)
        
        filename = f"image_{uuid.uuid4().hex[:8]}.png"
        filepath = GALLERY_DIR / filename
        image.save(str(filepath), quality=95)
        
        # Save to database
        save_generated_file('image', filename, f"/gallery/{filename}", prompt[:200])
        
        task_manager.update_task(task_id, step=4, step_name='Complete!', progress=100, status='completed', result=f"/gallery/{filename}")
        
        return f"/gallery/{filename}", None
        
    except InterruptedError:
        task_manager.update_task(task_id, status='cancelled')
        return None, "Cancelled"
    except Exception as e:
        print(f"[SD] Generation error: {e}")
        traceback.print_exc()
        task_manager.update_task(task_id, status='error', error=str(e))
        return None, str(e)

def generate_image(prompt, num_steps=None):
    """Generate image from prompt (legacy, no progress)."""
    import torch
    num_steps = num_steps or SD_STEPS
    
    pipeline = load_sd_pipeline()
    if pipeline is None:
        return None, "Image generation not available"
    
    try:
        with torch.inference_mode():
            result = pipeline(
                prompt=prompt,
                num_inference_steps=num_steps,
                guidance_scale=SD_GUIDANCE,
                height=1024,
                width=1024
            )
        image = result.images[0]
        filename = f"image_{uuid.uuid4().hex[:8]}.png"
        filepath = GALLERY_DIR / filename
        image.save(str(filepath))
        save_generated_file('image', filename, f"/gallery/{filename}", prompt[:200])
        return f"/gallery/{filename}", None
    except Exception as e:
        print(f"[SD] Generation error: {e}")
        return None, str(e)

# =============================================================================
# Video Generation (Stable Video Diffusion)
# =============================================================================

def generate_video_with_progress(prompt_or_image, task_id):
    """Generate video from prompt or image using Stable Video Diffusion."""
    import torch
    from PIL import Image
    
    task_manager.update_task(task_id, step=1, step_name='Loading video generation model...', progress=5)
    
    try:
        # First generate an image if we got a text prompt
        if isinstance(prompt_or_image, str):
            task_manager.update_task(task_id, step=1, step_name='Generating base image from prompt...', progress=10)
            
            sd_pipe = load_sd_pipeline()
            if sd_pipe is None:
                task_manager.update_task(task_id, status='error', error='Image generation not available')
                return None, "Image generation not available"
            
            if task_manager.is_cancelled(task_id):
                task_manager.update_task(task_id, status='cancelled')
                return None, "Cancelled"
            
            with torch.inference_mode():
                result = sd_pipe(
                    prompt=prompt_or_image,
                    num_inference_steps=20,  # Fewer steps for base image
                    guidance_scale=7.5,
                    height=576,  # SVD preferred size
                    width=1024
                )
            base_image = result.images[0]
        else:
            base_image = prompt_or_image
        
        if task_manager.is_cancelled(task_id):
            task_manager.update_task(task_id, status='cancelled')
            return None, "Cancelled"
        
        task_manager.update_task(task_id, step=2, step_name='Loading Stable Video Diffusion...', progress=30)
        
        video_pipe = load_video_pipeline()
        if video_pipe is None:
            task_manager.update_task(task_id, status='error', error='Video generation not available')
            return None, "Video generation not available"
        
        # Resize image for SVD
        base_image = base_image.resize((1024, 576))
        
        task_manager.update_task(task_id, step=3, step_name='Generating video frames...', progress=40)
        
        if task_manager.is_cancelled(task_id):
            task_manager.update_task(task_id, status='cancelled')
            return None, "Cancelled"
        
        # Generate video frames
        with torch.inference_mode():
            frames = video_pipe(
                base_image,
                num_frames=VIDEO_FRAMES,
                decode_chunk_size=8,  # Memory optimization
                motion_bucket_id=127,  # Motion amount
                noise_aug_strength=0.02
            ).frames[0]
        
        task_manager.update_task(task_id, step=4, step_name='Encoding video...', progress=80)
        
        if task_manager.is_cancelled(task_id):
            task_manager.update_task(task_id, status='cancelled')
            return None, "Cancelled"
        
        # Save video
        from diffusers.utils import export_to_video
        
        filename = f"video_{uuid.uuid4().hex[:8]}.mp4"
        filepath = VIDEO_DIR / filename
        export_to_video(frames, str(filepath), fps=VIDEO_FPS)
        
        # Save to database
        prompt_text = prompt_or_image if isinstance(prompt_or_image, str) else "image-to-video"
        save_generated_file('video', filename, f"/videos/{filename}", prompt_text[:200])
        
        task_manager.update_task(task_id, step=5, step_name='Complete!', progress=100, status='completed',
                                result={'url': f"/videos/{filename}", 'filename': filename})
        
        return f"/videos/{filename}", filename
        
    except InterruptedError:
        task_manager.update_task(task_id, status='cancelled')
        return None, "Cancelled"
    except Exception as e:
        print(f"[SVD] Video generation error: {e}")
        traceback.print_exc()
        task_manager.update_task(task_id, status='error', error=str(e))
        return None, str(e)

# =============================================================================
# Document Generation with Progress
# =============================================================================

def generate_pdf_with_progress(title, content, task_id):
    """Generate PDF document with progress tracking."""
    try:
        task_manager.update_task(task_id, step=1, step_name='Importing PDF libraries...', progress=10)
        
        if task_manager.is_cancelled(task_id):
            return None, "Cancelled"
        
        from reportlab.lib.pagesizes import letter
        from reportlab.pdfgen import canvas
        from reportlab.lib.units import inch
        
        task_manager.update_task(task_id, step=2, step_name='Creating document...', progress=30)
        
        filename = f"doc_{uuid.uuid4().hex[:8]}.pdf"
        filepath = DOCUMENT_DIR / filename
        
        if task_manager.is_cancelled(task_id):
            return None, "Cancelled"
        
        task_manager.update_task(task_id, step=3, step_name='Writing content...', progress=50)
        
        c = canvas.Canvas(str(filepath), pagesize=letter)
        width, height = letter
        
        c.setFont("Helvetica-Bold", 18)
        c.drawString(inch, height - inch, title)
        
        c.setFont("Helvetica", 12)
        y = height - 1.5 * inch
        lines = content.split('\n')
        total_lines = len(lines)
        
        for i, line in enumerate(lines):
            if task_manager.is_cancelled(task_id):
                return None, "Cancelled"
            
            if y < inch:
                c.showPage()
                y = height - inch
            c.drawString(inch, y, line[:80])
            y -= 14
            
            if i % 10 == 0:
                progress = 50 + int((i / max(total_lines, 1)) * 40)
                task_manager.update_task(task_id, step=3, step_name=f'Writing content... ({i}/{total_lines} lines)', progress=progress)
        
        task_manager.update_task(task_id, step=4, step_name='Saving file...', progress=95)
        c.save()
        
        task_manager.update_task(task_id, step=5, step_name='Complete!', progress=100, status='completed', result={'url': f"/documents/{filename}", 'filename': filename})
        
        return f"/documents/{filename}", filename
        
    except Exception as e:
        print(f"[PDF] Generation error: {e}")
        task_manager.update_task(task_id, status='error', error=str(e))
        return None, str(e)

def generate_pdf(title, content):
    """Generate PDF document (legacy)."""
    try:
        from reportlab.lib.pagesizes import letter
        from reportlab.pdfgen import canvas
        from reportlab.lib.units import inch
        
        filename = f"doc_{uuid.uuid4().hex[:8]}.pdf"
        filepath = DOCUMENT_DIR / filename
        
        c = canvas.Canvas(str(filepath), pagesize=letter)
        width, height = letter
        
        # Title
        c.setFont("Helvetica-Bold", 18)
        c.drawString(inch, height - inch, title)
        
        # Content
        c.setFont("Helvetica", 12)
        y = height - 1.5 * inch
        for line in content.split('\n'):
            if y < inch:
                c.showPage()
                y = height - inch
            c.drawString(inch, y, line[:80])
            y -= 14
        
        c.save()
        return f"/documents/{filename}", filename
    except Exception as e:
        print(f"[PDF] Generation error: {e}")
        return None, str(e)

def generate_docx(title, content):
    """Generate Word document."""
    try:
        from docx import Document
        
        filename = f"doc_{uuid.uuid4().hex[:8]}.docx"
        filepath = DOCUMENT_DIR / filename
        
        doc = Document()
        doc.add_heading(title, 0)
        
        for para in content.split('\n\n'):
            doc.add_paragraph(para)
        
        doc.save(str(filepath))
        return f"/documents/{filename}", filename
    except Exception as e:
        print(f"[DOCX] Generation error: {e}")
        return None, str(e)

def generate_pptx(title, slides_content):
    """Generate PowerPoint presentation."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        
        filename = f"presentation_{uuid.uuid4().hex[:8]}.pptx"
        filepath = DOCUMENT_DIR / filename
        
        prs = Presentation()
        
        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        slide.shapes.title.text = title
        
        # Content slides
        bullet_layout = prs.slide_layouts[1]
        for slide_content in slides_content:
            slide = prs.slides.add_slide(bullet_layout)
            slide.shapes.title.text = slide_content.get('title', 'Slide')
            body = slide.shapes.placeholders[1]
            tf = body.text_frame
            tf.text = slide_content.get('content', '')
        
        prs.save(str(filepath))
        return f"/documents/{filename}", filename
    except Exception as e:
        print(f"[PPTX] Generation error: {e}")
        return None, str(e)

# =============================================================================
# Audio Generation (Text-to-Speech)
# =============================================================================

AUDIO_DIR = WORKSPACE_DIR / "audio"
AUDIO_DIR.mkdir(parents=True, exist_ok=True)

def generate_audio(text, task_id=None, lang='en'):
    """Generate audio from text using gTTS."""
    try:
        if task_id:
            task_manager.update_task(task_id, step=1, step_name='Initializing text-to-speech...', progress=10)
        
        from gtts import gTTS
        
        if task_manager and task_id and task_manager.is_cancelled(task_id):
            return None, "Cancelled"
        
        if task_id:
            task_manager.update_task(task_id, step=2, step_name='Converting text to speech...', progress=30)
        
        # Create audio
        tts = gTTS(text=text, lang=lang, slow=False)
        
        if task_id:
            task_manager.update_task(task_id, step=3, step_name='Saving audio file...', progress=70)
        
        filename = f"audio_{uuid.uuid4().hex[:8]}.mp3"
        filepath = AUDIO_DIR / filename
        tts.save(str(filepath))
        
        # Save to database
        save_generated_file('audio', filename, f"/audio/{filename}", text[:200])
        
        if task_id:
            task_manager.update_task(task_id, step=4, step_name='Complete!', progress=100, status='completed', 
                                    result={'url': f"/audio/{filename}", 'filename': filename})
        
        return f"/audio/{filename}", filename
    except Exception as e:
        print(f"[Audio] Generation error: {e}")
        if task_id:
            task_manager.update_task(task_id, status='error', error=str(e))
        return None, str(e)

# =============================================================================
# QR Code Generation
# =============================================================================

def generate_qrcode(data, task_id=None):
    """Generate QR code image."""
    try:
        if task_id:
            task_manager.update_task(task_id, step=1, step_name='Creating QR code...', progress=30)
        
        import qrcode
        from qrcode.image.styledpil import StyledPilImage
        from qrcode.image.styles.moduledrawers import RoundedModuleDrawer
        from qrcode.image.styles.colormasks import RadialGradiantColorMask
        
        qr = qrcode.QRCode(version=1, error_correction=qrcode.constants.ERROR_CORRECT_H, box_size=10, border=4)
        qr.add_data(data)
        qr.make(fit=True)
        
        if task_id:
            task_manager.update_task(task_id, step=2, step_name='Styling QR code...', progress=60)
        
        # Create styled QR code with gradient
        try:
            img = qr.make_image(
                image_factory=StyledPilImage,
                module_drawer=RoundedModuleDrawer(),
                color_mask=RadialGradiantColorMask(
                    back_color=(255, 255, 255),
                    center_color=(139, 92, 246),
                    edge_color=(236, 72, 153)
                )
            )
        except:
            # Fallback to simple QR
            img = qr.make_image(fill_color="#8b5cf6", back_color="white")
        
        if task_id:
            task_manager.update_task(task_id, step=3, step_name='Saving image...', progress=90)
        
        filename = f"qr_{uuid.uuid4().hex[:8]}.png"
        filepath = GALLERY_DIR / filename
        img.save(str(filepath))
        
        save_generated_file('qr', filename, f"/gallery/{filename}", data[:200])
        
        if task_id:
            task_manager.update_task(task_id, step=4, step_name='Complete!', progress=100, status='completed',
                                    result={'url': f"/gallery/{filename}", 'filename': filename})
        
        return f"/gallery/{filename}", filename
    except Exception as e:
        print(f"[QR] Generation error: {e}")
        if task_id:
            task_manager.update_task(task_id, status='error', error=str(e))
        return None, str(e)

# =============================================================================
# Chart Generation
# =============================================================================

def generate_chart(chart_type, data, title="Chart", task_id=None):
    """Generate chart/graph image."""
    try:
        if task_id:
            task_manager.update_task(task_id, step=1, step_name='Setting up chart...', progress=20)
        
        import matplotlib
        matplotlib.use('Agg')  # Non-interactive backend
        import matplotlib.pyplot as plt
        import numpy as np
        
        # Set dark theme
        plt.style.use('dark_background')
        fig, ax = plt.subplots(figsize=(10, 6), facecolor='#12121a')
        ax.set_facecolor('#1a1a25')
        
        if task_id:
            task_manager.update_task(task_id, step=2, step_name='Creating visualization...', progress=50)
        
        # Parse data if it's a string
        if isinstance(data, str):
            try:
                data = json.loads(data)
            except:
                # Try to parse as simple format: "label1:value1,label2:value2"
                data = {}
                for item in data.split(','):
                    if ':' in item:
                        k, v = item.split(':', 1)
                        try:
                            data[k.strip()] = float(v.strip())
                        except:
                            data[k.strip()] = v.strip()
        
        labels = list(data.keys()) if isinstance(data, dict) else [f"Item {i+1}" for i in range(len(data))]
        values = list(data.values()) if isinstance(data, dict) else data
        
        # Gradient colors
        colors = plt.cm.plasma(np.linspace(0.2, 0.8, len(labels)))
        
        if chart_type == 'bar':
            bars = ax.bar(labels, values, color=colors, edgecolor='white', linewidth=0.5)
            ax.set_xlabel('Categories', color='#a0a0b0')
            ax.set_ylabel('Values', color='#a0a0b0')
        elif chart_type == 'pie':
            wedges, texts, autotexts = ax.pie(values, labels=labels, autopct='%1.1f%%', colors=colors,
                                              textprops={'color': 'white'}, wedgeprops={'edgecolor': '#12121a'})
        elif chart_type == 'line':
            ax.plot(labels, values, color='#8b5cf6', linewidth=2, marker='o', markersize=8,
                   markerfacecolor='#ec4899', markeredgecolor='white')
            ax.fill_between(labels, values, alpha=0.3, color='#8b5cf6')
            ax.set_xlabel('X', color='#a0a0b0')
            ax.set_ylabel('Y', color='#a0a0b0')
        else:
            ax.bar(labels, values, color=colors)
        
        ax.set_title(title, color='white', fontsize=14, fontweight='bold', pad=20)
        ax.tick_params(colors='#606070')
        for spine in ax.spines.values():
            spine.set_color('#2a2a35')
        
        plt.tight_layout()
        
        if task_id:
            task_manager.update_task(task_id, step=3, step_name='Saving chart...', progress=85)
        
        filename = f"chart_{uuid.uuid4().hex[:8]}.png"
        filepath = GALLERY_DIR / filename
        plt.savefig(str(filepath), dpi=150, facecolor='#12121a', edgecolor='none', bbox_inches='tight')
        plt.close()
        
        save_generated_file('chart', filename, f"/gallery/{filename}", f"{chart_type}: {title}")
        
        if task_id:
            task_manager.update_task(task_id, step=4, step_name='Complete!', progress=100, status='completed',
                                    result={'url': f"/gallery/{filename}", 'filename': filename})
        
        return f"/gallery/{filename}", filename
    except Exception as e:
        print(f"[Chart] Generation error: {e}")
        if task_id:
            task_manager.update_task(task_id, status='error', error=str(e))
        return None, str(e)

# =============================================================================
# Excel Generation
# =============================================================================

def generate_excel(data, title="Spreadsheet", task_id=None):
    """Generate Excel spreadsheet."""
    try:
        if task_id:
            task_manager.update_task(task_id, step=1, step_name='Creating spreadsheet...', progress=30)
        
        import pandas as pd
        
        # Parse data
        if isinstance(data, str):
            try:
                data = json.loads(data)
            except:
                # Create from text (each line is a row, comma-separated)
                rows = [line.split(',') for line in data.strip().split('\n')]
                if rows:
                    headers = rows[0] if len(rows) > 1 else [f"Col{i+1}" for i in range(len(rows[0]))]
                    data = {headers[i]: [row[i] if i < len(row) else '' for row in rows[1:]] for i in range(len(headers))}
        
        df = pd.DataFrame(data)
        
        if task_id:
            task_manager.update_task(task_id, step=2, step_name='Formatting...', progress=60)
        
        filename = f"spreadsheet_{uuid.uuid4().hex[:8]}.xlsx"
        filepath = DOCUMENT_DIR / filename
        
        # Save with formatting
        with pd.ExcelWriter(str(filepath), engine='xlsxwriter') as writer:
            df.to_excel(writer, sheet_name='Data', index=False)
            workbook = writer.book
            worksheet = writer.sheets['Data']
            
            # Format header
            header_format = workbook.add_format({
                'bold': True, 'bg_color': '#8b5cf6', 'font_color': 'white',
                'border': 1, 'align': 'center'
            })
            for col_num, col_name in enumerate(df.columns):
                worksheet.write(0, col_num, col_name, header_format)
                worksheet.set_column(col_num, col_num, max(len(str(col_name)), 12))
        
        save_generated_file('xlsx', filename, f"/documents/{filename}", title)
        
        if task_id:
            task_manager.update_task(task_id, step=3, step_name='Complete!', progress=100, status='completed',
                                    result={'url': f"/documents/{filename}", 'filename': filename})
        
        return f"/documents/{filename}", filename
    except Exception as e:
        print(f"[Excel] Generation error: {e}")
        if task_id:
            task_manager.update_task(task_id, status='error', error=str(e))
        return None, str(e)

# =============================================================================
# API Routes
# =============================================================================

@app.route('/')
def index():
    return render_template_string(HTML_TEMPLATE)

@app.route('/api/providers', methods=['GET'])
def get_providers():
    """Get all providers with their models and configured status."""
    api_keys = get_all_api_keys()
    custom = get_custom_models()
    
    result = {}
    for pid, pdata in PROVIDERS.items():
        models = list(pdata['models'])
        for cm in custom:
            if cm['provider'] == pid:
                models.append(cm['model_id'])
        result[pid] = {
            **pdata,
            'models': models,
            'configured': pid in api_keys and bool(api_keys[pid].get('api_key'))
        }
    return jsonify(result)

@app.route('/api/ollama/models', methods=['GET'])
def get_ollama_models():
    """Get list of Ollama models with their download status."""
    try:
        url = OLLAMA_HOST + "/api/tags"
        req = urllib.request.Request(url, method='GET')
        with urllib.request.urlopen(req, timeout=5) as response:
            data = json.loads(response.read().decode('utf-8'))
            downloaded = [m['name'] for m in data.get('models', [])]
            return jsonify({
                'available': True,
                'downloaded': downloaded
            })
    except Exception as e:
        return jsonify({
            'available': False,
            'downloaded': [],
            'error': str(e)
        })

@app.route('/api/ollama/pull', methods=['POST'])
def pull_ollama_model():
    """Start pulling an Ollama model - returns immediately, client streams progress."""
    data = request.json
    model = data.get('model')
    if not model:
        return jsonify({'error': 'No model specified'}), 400
    return jsonify({'status': 'started', 'model': model})

@app.route('/api/ollama/pull/stream/<model>')
def stream_ollama_pull(model):
    """Stream Ollama model pull progress."""
    def generate():
        try:
            url = OLLAMA_HOST + "/api/pull"
            payload = json.dumps({"name": model, "stream": True}).encode('utf-8')
            req = urllib.request.Request(url, payload, {"Content-Type": "application/json"}, method='POST')
            
            with urllib.request.urlopen(req, timeout=600) as response:
                buffer = ""
                for chunk in iter(lambda: response.read(1024), b''):
                    buffer += chunk.decode('utf-8')
                    while '\n' in buffer:
                        line, buffer = buffer.split('\n', 1)
                        if line.strip():
                            try:
                                data = json.loads(line)
                                yield f"data: {json.dumps(data)}\n\n"
                            except:
                                pass
            yield f"data: {json.dumps({'status': 'success'})}\n\n"
        except Exception as e:
            yield f"data: {json.dumps({'error': str(e)})}\n\n"
    
    return Response(generate(), mimetype='text/event-stream')

@app.route('/api/settings', methods=['GET'])
def get_settings():
    """Get all settings including API keys (masked) and current model."""
    api_keys = get_all_api_keys()
    masked = {}
    for p, data in api_keys.items():
        key = data.get('api_key', '')
        masked[p] = {'configured': bool(key), 'masked_key': f"...{key[-4:]}" if len(key) > 4 else "****", 'base_url': data.get('base_url', '')}
    
    return jsonify({
        'api_keys': masked,
        'current_provider': get_setting('current_provider', 'openai'),
        'current_model': get_setting('current_model', 'gpt-4o-mini'),
        'custom_models': get_custom_models()
    })

@app.route('/api/settings/apikey', methods=['POST'])
def save_api_key_route():
    """Save an API key for a provider."""
    data = request.json
    provider = data.get('provider')
    api_key = data.get('api_key')
    base_url = data.get('base_url')
    
    if not provider or not api_key:
        return jsonify({'error': 'Provider and API key required'}), 400
    
    save_api_key(provider, api_key, base_url)
    return jsonify({'success': True})

@app.route('/api/settings/apikey/<provider>', methods=['DELETE'])
def delete_api_key_route(provider):
    """Delete an API key."""
    conn = get_db()
    conn.execute('DELETE FROM api_keys WHERE provider = ?', (provider,))
    conn.commit()
    conn.close()
    return jsonify({'success': True})

@app.route('/api/settings/model', methods=['POST'])
def set_current_model():
    """Set the current provider and model."""
    data = request.json
    provider = data.get('provider')
    model = data.get('model')
    
    if provider: set_setting('current_provider', provider)
    if model: set_setting('current_model', model)
    
    return jsonify({'success': True})

@app.route('/api/models/custom', methods=['POST'])
def add_custom_model_route():
    """Add a custom model."""
    data = request.json
    add_custom_model(data['provider'], data['model_id'], data.get('display_name'))
    return jsonify({'success': True})

@app.route('/api/models/custom/<model_id>', methods=['DELETE'])
def delete_custom_model_route(model_id):
    """Delete a custom model."""
    delete_custom_model(model_id)
    return jsonify({'success': True})

@app.route('/api/test/<provider>', methods=['POST'])
def test_provider(provider):
    """Test API connection for a provider."""
    creds = get_api_key(provider)
    if not creds and provider != 'ollama':
        return jsonify({'success': False, 'error': 'No API key configured'})
    
    api_key = creds['api_key'] if creds else None
    base_url = creds.get('base_url') if creds else None
    
    try:
        test_messages = [{"role": "user", "content": "Say 'OK' and nothing else."}]
        
        if provider == "openai":
            url = base_url or "https://api.openai.com/v1/chat/completions"
            payload = {"model": "gpt-3.5-turbo", "messages": test_messages, "max_tokens": 10}
            headers = {"Content-Type": "application/json", "Authorization": f"Bearer {api_key}"}
            req = urllib.request.Request(url, json.dumps(payload).encode('utf-8'), headers, method='POST')
            with urllib.request.urlopen(req, timeout=30) as resp:
                data = json.loads(resp.read().decode('utf-8'))
                model_used = data.get('model', 'unknown')
                return jsonify({'success': True, 'message': f'Connected! Model: {model_used}'})
        
        elif provider == "anthropic":
            url = "https://api.anthropic.com/v1/messages"
            payload = {"model": "claude-3-haiku-20240307", "max_tokens": 10, "messages": test_messages}
            headers = {"Content-Type": "application/json", "x-api-key": api_key, "anthropic-version": "2023-06-01"}
            req = urllib.request.Request(url, json.dumps(payload).encode('utf-8'), headers, method='POST')
            with urllib.request.urlopen(req, timeout=30) as resp:
                data = json.loads(resp.read().decode('utf-8'))
                model_used = data.get('model', 'unknown')
                return jsonify({'success': True, 'message': f'Connected! Model: {model_used}'})
        
        elif provider == "google":
            url = f"https://generativelanguage.googleapis.com/v1beta/models/gemini-1.5-flash:generateContent?key={api_key}"
            payload = {"contents": [{"parts": [{"text": "Say OK"}]}]}
            headers = {"Content-Type": "application/json"}
            req = urllib.request.Request(url, json.dumps(payload).encode('utf-8'), headers, method='POST')
            with urllib.request.urlopen(req, timeout=30) as resp:
                data = json.loads(resp.read().decode('utf-8'))
                return jsonify({'success': True, 'message': 'Connected to Gemini API!'})
        
        elif provider == "xai":
            url = "https://api.x.ai/v1/chat/completions"
            payload = {"model": "grok-beta", "messages": test_messages, "max_tokens": 10}
            headers = {"Content-Type": "application/json", "Authorization": f"Bearer {api_key}"}
            req = urllib.request.Request(url, json.dumps(payload).encode('utf-8'), headers, method='POST')
            with urllib.request.urlopen(req, timeout=30) as resp:
                data = json.loads(resp.read().decode('utf-8'))
                return jsonify({'success': True, 'message': 'Connected to Grok API!'})
        
        elif provider == "ollama":
            url = (base_url or "http://localhost:11434") + "/api/tags"
            req = urllib.request.Request(url, method='GET')
            with urllib.request.urlopen(req, timeout=10) as resp:
                data = json.loads(resp.read().decode('utf-8'))
                models = [m['name'] for m in data.get('models', [])]
                count = len(models)
                return jsonify({'success': True, 'message': f'Connected! {count} models available', 'models': models[:5]})
        
        return jsonify({'success': False, 'error': 'Unknown provider'})
        
    except urllib.error.HTTPError as e:
        error_body = ""
        try:
            error_body = e.read().decode('utf-8')
            error_data = json.loads(error_body)
            error_msg = error_data.get('error', {}).get('message', str(e.reason))
        except:
            error_msg = f"{e.code}: {e.reason}"
        return jsonify({'success': False, 'error': error_msg})
    except urllib.error.URLError as e:
        return jsonify({'success': False, 'error': f'Connection failed: {str(e.reason)}'})
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)})

@app.route('/api/task/<task_id>', methods=['GET'])
def get_task_status(task_id):
    """Get status of a running task."""
    task = task_manager.get_task(task_id)
    if not task:
        return jsonify({'error': 'Task not found'}), 404
    return jsonify(task)

@app.route('/api/task/<task_id>/cancel', methods=['POST'])
def cancel_task(task_id):
    """Cancel a running task."""
    if task_manager.cancel_task(task_id):
        return jsonify({'success': True, 'message': 'Cancellation requested'})
    return jsonify({'error': 'Task not found'}), 404

@app.route('/api/task/stream/<task_id>')
def stream_task_progress(task_id):
    """Stream task progress via SSE."""
    def generate():
        last_status = None
        while True:
            task = task_manager.get_task(task_id)
            if not task:
                yield f"data: {json.dumps({'error': 'Task not found'})}\n\n"
                break
            
            # Only send updates when something changed
            current_status = (task.get('status'), task.get('step'), task.get('progress'))
            if current_status != last_status:
                yield f"data: {json.dumps(task)}\n\n"
                last_status = current_status
            
            # Exit if task is done
            if task.get('status') in ['completed', 'error', 'cancelled']:
                break
            
            time.sleep(0.2)  # Poll every 200ms
    
    return Response(generate(), mimetype='text/event-stream')

@app.route('/api/conversations', methods=['GET'])
def list_conversations():
    conn = get_db()
    rows = conn.execute('SELECT * FROM conversations ORDER BY updated_at DESC').fetchall()
    conn.close()
    return jsonify([dict(row) for row in rows])

@app.route('/api/conversation/create', methods=['POST'])
def create_conversation():
    conversation_id = str(uuid.uuid4())
    now = datetime.now().isoformat()
    conn = get_db()
    conn.execute('INSERT INTO conversations (id, title, created_at, updated_at, message_count) VALUES (?, ?, ?, ?, ?)',
                 (conversation_id, 'New Conversation', now, now, 0))
    conn.commit()
    conn.close()
    return jsonify({'id': conversation_id, 'title': 'New Conversation'})

@app.route('/api/conversation/<conversation_id>', methods=['GET'])
def get_conversation(conversation_id):
    conn = get_db()
    rows = conn.execute('SELECT * FROM messages WHERE conversation_id = ? ORDER BY created_at', (conversation_id,)).fetchall()
    messages = []
    for row in rows:
        msg = dict(row)
        if msg['attachments']: msg['attachments'] = json.loads(msg['attachments'])
        messages.append(msg)
    conn.close()
    return jsonify(messages)

@app.route('/api/conversation/<conversation_id>', methods=['DELETE'])
def delete_conversation(conversation_id):
    conn = get_db()
    conn.execute('DELETE FROM messages WHERE conversation_id = ?', (conversation_id,))
    conn.execute('DELETE FROM conversations WHERE id = ?', (conversation_id,))
    conn.commit()
    conn.close()
    return jsonify({'success': True})

@app.route('/api/chat/stream', methods=['POST'])
def stream_chat():
    data = request.json
    message = data.get('message', '')
    conversation_id = data.get('conversation_id')
    req_attachments = data.get('attachments', [])
    req_provider = data.get('provider')
    req_model = data.get('model')
    image_gen_enabled = data.get('image_gen_enabled', True)
    video_gen_enabled = data.get('video_gen_enabled', True)
    
    if not conversation_id:
        return jsonify({'error': 'No conversation ID'}), 400
    
    # Use request-specified or saved settings
    provider = req_provider or get_setting('current_provider', 'openai')
    model = req_model or get_setting('current_model', 'gemini-2.5-pro')
    
    def generate():
        conn = get_db()
        now = datetime.now().isoformat()
        
        conn.execute('INSERT INTO messages (conversation_id, role, content, attachments, created_at) VALUES (?, ?, ?, ?, ?)',
                     (conversation_id, 'user', message, json.dumps(req_attachments), now))
        conn.commit()
        
        cursor = conn.execute('SELECT role, content FROM messages WHERE conversation_id = ? ORDER BY created_at DESC LIMIT 10', (conversation_id,))
        history = list(reversed([dict(row) for row in cursor.fetchall()]))
        
        messages = [{"role": "system", "content": SYSTEM_PROMPT}]
        for msg in history:
            messages.append({"role": msg['role'], "content": msg['content']})
        
        lower_msg = message.lower()
        
        # Image generation with progress tracking
        if any(phrase in lower_msg for phrase in ['generate image', 'create image', 'draw', 'make an image']):
            if not image_gen_enabled:
                yield f"data: {json.dumps({'content': '‚ö†Ô∏è Image generation is currently disabled. Enable it using the üé® toggle in the header.'})}\n\n"
                yield f"data: {json.dumps({'done': True})}\n\n"
                return
            if ENABLE_IMAGE_GEN:
                task_id = task_manager.create_task('image_generation')
                yield f"data: {json.dumps({'task_start': {'id': task_id, 'type': 'image', 'title': 'üé® Image Generation', 'steps': ['Loading Stable Diffusion model', 'Generating image', 'Post-processing', 'Saving to gallery']}})}\n\n"
                
                # Run generation in thread and stream progress
                def run_generation():
                    generate_image_with_progress(message, task_id)
                
                thread = threading.Thread(target=run_generation)
                thread.start()
                
                # Stream progress updates
                last_progress = -1
                while thread.is_alive() or task_manager.get_task(task_id).get('status') == 'running':
                    task = task_manager.get_task(task_id)
                    if task.get('progress', 0) != last_progress:
                        yield f"data: {json.dumps({'task_progress': task})}\n\n"
                        last_progress = task.get('progress', 0)
                    time.sleep(0.1)
                
                thread.join()
                
                # Final result
                task = task_manager.get_task(task_id)
                if task.get('status') == 'completed' and task.get('result'):
                    yield f"data: {json.dumps({'image': task['result']})}\n\n"
                    yield f"data: {json.dumps({'content': '‚úÖ Image generated successfully!'})}\n\n"
                elif task.get('status') == 'cancelled':
                    yield f"data: {json.dumps({'content': '‚ö†Ô∏è Image generation was cancelled.'})}\n\n"
                else:
                    error_msg = task.get('error', 'Unknown error')
                    yield f"data: {json.dumps({'content': f'‚ùå Failed: {error_msg}'})}\n\n"
                
                yield f"data: {json.dumps({'task_complete': task_id})}\n\n"
            else:
                yield f"data: {json.dumps({'content': '‚ö†Ô∏è Image generation is not enabled. Set `ENABLE_IMAGE_GENERATION=true` in your .env file and install the required dependencies (torch, diffusers, transformers).'})}\n\n"
            yield f"data: {json.dumps({'done': True})}\n\n"
            return
        
        # Video generation with progress tracking
        if any(phrase in lower_msg for phrase in ['generate video', 'create video', 'make a video', 'animate']):
            if not video_gen_enabled:
                yield f"data: {json.dumps({'content': '‚ö†Ô∏è Video generation is currently disabled. Enable it using the üé¨ toggle in the header.'})}\n\n"
                yield f"data: {json.dumps({'done': True})}\n\n"
                return
            if ENABLE_VIDEO_GEN:
                task_id = task_manager.create_task('video_generation')
                yield f"data: {json.dumps({'task_start': {'id': task_id, 'type': 'video', 'title': 'üé¨ Video Generation', 'steps': ['Generating base image', 'Loading video model', 'Generating frames', 'Encoding video', 'Saving']}})}\n\n"
                
                def run_video():
                    generate_video_with_progress(message, task_id)
                
                thread = threading.Thread(target=run_video)
                thread.start()
                
                last_progress = -1
                while thread.is_alive() or task_manager.get_task(task_id).get('status') == 'running':
                    task = task_manager.get_task(task_id)
                    if task.get('progress', 0) != last_progress:
                        yield f"data: {json.dumps({'task_progress': task})}\n\n"
                        last_progress = task.get('progress', 0)
                    time.sleep(0.2)
                
                thread.join()
                
                task = task_manager.get_task(task_id)
                if task.get('status') == 'completed' and task.get('result'):
                    yield f"data: {json.dumps({'video': task['result']['url']})}\n\n"
                    yield f"data: {json.dumps({'content': '‚úÖ Video generated successfully! Click to play.'})}\n\n"
                elif task.get('status') == 'cancelled':
                    yield f"data: {json.dumps({'content': '‚ö†Ô∏è Video generation was cancelled.'})}\n\n"
                else:
                    error_msg = task.get('error', 'Unknown error')
                    yield f"data: {json.dumps({'content': f'‚ùå Failed: {error_msg}'})}\n\n"
                
                yield f"data: {json.dumps({'task_complete': task_id})}\n\n"
            else:
                yield f"data: {json.dumps({'content': '‚ö†Ô∏è Video generation is not enabled. Set `ENABLE_VIDEO_GENERATION=true` in your .env file.'})}\n\n"
            yield f"data: {json.dumps({'done': True})}\n\n"
            return
        
        # PDF generation with progress tracking
        if any(phrase in lower_msg for phrase in ['generate pdf', 'create pdf']):
            task_id = task_manager.create_task('pdf_generation')
            yield f"data: {json.dumps({'task_start': {'id': task_id, 'type': 'document', 'title': 'üìÑ PDF Generation', 'steps': ['Importing libraries', 'Creating document', 'Writing content', 'Saving file']}})}\n\n"
            
            def run_pdf():
                generate_pdf_with_progress("Generated Document", message, task_id)
            
            thread = threading.Thread(target=run_pdf)
            thread.start()
            
            last_progress = -1
            while thread.is_alive():
                task = task_manager.get_task(task_id)
                if task.get('progress', 0) != last_progress:
                    yield f"data: {json.dumps({'task_progress': task})}\n\n"
                    last_progress = task.get('progress', 0)
                time.sleep(0.1)
            
            thread.join()
            task = task_manager.get_task(task_id)
            
            if task.get('status') == 'completed' and task.get('result'):
                yield f"data: {json.dumps({'document': task['result']['url'], 'document_name': task['result']['filename']})}\n\n"
                yield f"data: {json.dumps({'content': '‚úÖ PDF document generated!'})}\n\n"
            elif task.get('status') == 'cancelled':
                yield f"data: {json.dumps({'content': '‚ö†Ô∏è PDF generation was cancelled.'})}\n\n"
            else:
                error_msg = task.get('error', 'Unknown error')
                yield f"data: {json.dumps({'content': f'‚ùå Failed: {error_msg}'})}\n\n"
            
            yield f"data: {json.dumps({'task_complete': task_id})}\n\n"
            yield f"data: {json.dumps({'done': True})}\n\n"
            return
        
        # Word document generation
        if any(phrase in lower_msg for phrase in ['generate word', 'create docx', 'create word']):
            task_id = task_manager.create_task('docx_generation')
            yield f"data: {json.dumps({'task_start': {'id': task_id, 'type': 'document', 'title': 'üìù Word Document Generation', 'steps': ['Importing libraries', 'Creating document', 'Writing content', 'Saving file']}})}\n\n"
            
            # For now, use sync version with simple progress
            task_manager.update_task(task_id, step=1, step_name='Creating document...', progress=25)
            yield f"data: {json.dumps({'task_progress': task_manager.get_task(task_id)})}\n\n"
            
            doc_url, filename = generate_docx("Generated Document", message)
            
            if doc_url:
                task_manager.update_task(task_id, step=4, step_name='Complete!', progress=100, status='completed')
                yield f"data: {json.dumps({'task_progress': task_manager.get_task(task_id)})}\n\n"
                yield f"data: {json.dumps({'document': doc_url, 'document_name': filename})}\n\n"
                yield f"data: {json.dumps({'content': '‚úÖ Word document generated!'})}\n\n"
            else:
                task_manager.update_task(task_id, status='error', error=filename)
                yield f"data: {json.dumps({'content': f'‚ùå Failed: {filename}'})}\n\n"
            
            yield f"data: {json.dumps({'task_complete': task_id})}\n\n"
            yield f"data: {json.dumps({'done': True})}\n\n"
            return
        
        # Audio/Speech generation
        if any(phrase in lower_msg for phrase in ['generate audio', 'create audio', 'text to speech', 'read aloud', 'speak this', 'say this']):
            task_id = task_manager.create_task('audio_generation')
            yield f"data: {json.dumps({'task_start': {'id': task_id, 'type': 'audio', 'title': 'üîä Audio Generation', 'steps': ['Initializing TTS', 'Converting text', 'Generating audio', 'Saving file']}})}\n\n"
            
            # Extract text to convert (remove command phrases)
            text_to_speak = message
            for phrase in ['generate audio', 'create audio', 'text to speech', 'read aloud', 'speak this', 'say this', 'for', 'of', ':']:
                text_to_speak = text_to_speak.lower().replace(phrase, '').strip()
            if not text_to_speak or len(text_to_speak) < 5:
                text_to_speak = message
            
            def run_audio():
                generate_audio(text_to_speak, task_id)
            
            thread = threading.Thread(target=run_audio)
            thread.start()
            
            last_progress = -1
            while thread.is_alive():
                task = task_manager.get_task(task_id)
                if task.get('progress', 0) != last_progress:
                    yield f"data: {json.dumps({'task_progress': task})}\n\n"
                    last_progress = task.get('progress', 0)
                time.sleep(0.1)
            
            thread.join()
            task = task_manager.get_task(task_id)
            
            if task.get('status') == 'completed' and task.get('result'):
                yield f"data: {json.dumps({'audio': task['result']['url'], 'audio_name': task['result']['filename']})}\n\n"
                yield f"data: {json.dumps({'content': '‚úÖ Audio generated! Click play to listen.'})}\n\n"
            else:
                error_msg = task.get('error', 'Unknown error')
                yield f"data: {json.dumps({'content': f'‚ùå Failed: {error_msg}'})}\n\n"
            
            yield f"data: {json.dumps({'task_complete': task_id})}\n\n"
            yield f"data: {json.dumps({'done': True})}\n\n"
            return
        
        # QR Code generation
        if any(phrase in lower_msg for phrase in ['generate qr', 'create qr', 'make qr', 'qr code']):
            task_id = task_manager.create_task('qr_generation')
            yield f"data: {json.dumps({'task_start': {'id': task_id, 'type': 'image', 'title': 'üì± QR Code Generation', 'steps': ['Creating QR code', 'Styling', 'Saving image']}})}\n\n"
            
            # Extract data for QR code
            qr_data = message
            for phrase in ['generate qr', 'create qr', 'make qr', 'qr code', 'for', 'of', 'with', ':']:
                qr_data = qr_data.lower().replace(phrase, '').strip()
            if not qr_data:
                qr_data = "https://github.com/LA-Rich/UltimateChat"
            
            def run_qr():
                generate_qrcode(qr_data, task_id)
            
            thread = threading.Thread(target=run_qr)
            thread.start()
            
            while thread.is_alive():
                task = task_manager.get_task(task_id)
                yield f"data: {json.dumps({'task_progress': task})}\n\n"
                time.sleep(0.1)
            
            thread.join()
            task = task_manager.get_task(task_id)
            
            if task.get('status') == 'completed' and task.get('result'):
                yield f"data: {json.dumps({'image': task['result']['url']})}\n\n"
                yield f"data: {json.dumps({'content': '‚úÖ QR code generated!'})}\n\n"
            else:
                error_msg = task.get('error', 'Unknown error')
                yield f"data: {json.dumps({'content': f'‚ùå Failed: {error_msg}'})}\n\n"
            
            yield f"data: {json.dumps({'task_complete': task_id})}\n\n"
            yield f"data: {json.dumps({'done': True})}\n\n"
            return
        
        # Chart generation
        if any(phrase in lower_msg for phrase in ['generate chart', 'create chart', 'make chart', 'create graph', 'bar chart', 'pie chart', 'line chart']):
            task_id = task_manager.create_task('chart_generation')
            yield f"data: {json.dumps({'task_start': {'id': task_id, 'type': 'image', 'title': 'üìä Chart Generation', 'steps': ['Setting up', 'Creating visualization', 'Rendering', 'Saving']}})}\n\n"
            
            # Determine chart type
            chart_type = 'bar'
            if 'pie' in lower_msg:
                chart_type = 'pie'
            elif 'line' in lower_msg:
                chart_type = 'line'
            
            # Try to extract data
            chart_data = {'Sample A': 30, 'Sample B': 45, 'Sample C': 25, 'Sample D': 60, 'Sample E': 35}
            chart_title = 'Generated Chart'
            
            # Look for data in format "label:value,label:value"
            import re
            data_match = re.search(r'data[:\s]+([^\n]+)', message, re.IGNORECASE)
            if data_match:
                try:
                    data_str = data_match.group(1)
                    chart_data = {}
                    for item in data_str.split(','):
                        if ':' in item:
                            k, v = item.split(':', 1)
                            chart_data[k.strip()] = float(v.strip())
                except:
                    pass
            
            title_match = re.search(r'title[:\s]+([^\n,]+)', message, re.IGNORECASE)
            if title_match:
                chart_title = title_match.group(1).strip()
            
            def run_chart():
                generate_chart(chart_type, chart_data, chart_title, task_id)
            
            thread = threading.Thread(target=run_chart)
            thread.start()
            
            while thread.is_alive():
                task = task_manager.get_task(task_id)
                yield f"data: {json.dumps({'task_progress': task})}\n\n"
                time.sleep(0.1)
            
            thread.join()
            task = task_manager.get_task(task_id)
            
            if task.get('status') == 'completed' and task.get('result'):
                yield f"data: {json.dumps({'image': task['result']['url']})}\n\n"
                yield f"data: {json.dumps({'content': f'‚úÖ {chart_type.title()} chart generated!'})}\n\n"
            else:
                error_msg = task.get('error', 'Unknown error')
                yield f"data: {json.dumps({'content': f'‚ùå Failed: {error_msg}'})}\n\n"
            
            yield f"data: {json.dumps({'task_complete': task_id})}\n\n"
            yield f"data: {json.dumps({'done': True})}\n\n"
            return
        
        # Excel generation
        if any(phrase in lower_msg for phrase in ['generate excel', 'create excel', 'create spreadsheet', 'make spreadsheet', 'generate xlsx']):
            task_id = task_manager.create_task('excel_generation')
            yield f"data: {json.dumps({'task_start': {'id': task_id, 'type': 'document', 'title': 'üìä Excel Generation', 'steps': ['Creating spreadsheet', 'Formatting', 'Saving']}})}\n\n"
            
            # Default sample data
            excel_data = {
                'Name': ['Item 1', 'Item 2', 'Item 3', 'Item 4'],
                'Value': [100, 200, 150, 300],
                'Category': ['A', 'B', 'A', 'C']
            }
            
            def run_excel():
                generate_excel(excel_data, "Generated Spreadsheet", task_id)
            
            thread = threading.Thread(target=run_excel)
            thread.start()
            
            while thread.is_alive():
                task = task_manager.get_task(task_id)
                yield f"data: {json.dumps({'task_progress': task})}\n\n"
                time.sleep(0.1)
            
            thread.join()
            task = task_manager.get_task(task_id)
            
            if task.get('status') == 'completed' and task.get('result'):
                yield f"data: {json.dumps({'document': task['result']['url'], 'document_name': task['result']['filename']})}\n\n"
                yield f"data: {json.dumps({'content': '‚úÖ Excel spreadsheet generated!'})}\n\n"
            else:
                error_msg = task.get('error', 'Unknown error')
                yield f"data: {json.dumps({'content': f'‚ùå Failed: {error_msg}'})}\n\n"
            
            yield f"data: {json.dumps({'task_complete': task_id})}\n\n"
            yield f"data: {json.dumps({'done': True})}\n\n"
            return
        
        # Stream LLM response using selected provider
        yield f"data: {json.dumps({'model_info': {'provider': provider, 'model': model}})}\n\n"
        
        full_response = ""
        for chunk in stream_response(provider, model, messages):
            full_response += chunk
            yield f"data: {json.dumps({'content': chunk})}\n\n"
        
        now = datetime.now().isoformat()
        conn.execute('INSERT INTO messages (conversation_id, role, content, attachments, created_at) VALUES (?, ?, ?, ?, ?)',
                     (conversation_id, 'assistant', full_response, '[]', now))
        
        message_count = conn.execute('SELECT COUNT(*) FROM messages WHERE conversation_id = ?', (conversation_id,)).fetchone()[0]
        
        # Smart title generation from first message
        title = generate_smart_title(message)
        conn.execute('UPDATE conversations SET title = ?, updated_at = ?, message_count = ? WHERE id = ? AND title = ?',
                     (title, now, message_count, conversation_id, 'New Conversation'))
        conn.commit()
        conn.close()
        
        yield f"data: {json.dumps({'done': True})}\n\n"
    
    return Response(generate(), mimetype='text/event-stream')

@app.route('/api/upload', methods=['POST'])
def upload_file():
    """Handle file upload."""
    if 'file' not in request.files:
        return jsonify({'error': 'No file provided'}), 400
    
    file = request.files['file']
    if file.filename == '':
        return jsonify({'error': 'No file selected'}), 400
    
    # Generate unique filename
    ext = Path(file.filename).suffix
    filename = f"{uuid.uuid4().hex[:12]}{ext}"
    filepath = UPLOAD_DIR / filename
    
    file.save(str(filepath))
    
    return jsonify({
        'url': f'/uploads/{filename}',
        'filename': file.filename,
        'size': filepath.stat().st_size
    })

@app.route('/uploads/<filename>')
def serve_upload(filename):
    """Serve uploaded files."""
    return send_from_directory(str(UPLOAD_DIR), filename)

@app.route('/gallery/<filename>')
def serve_gallery(filename):
    """Serve generated images."""
    return send_from_directory(str(GALLERY_DIR), filename)

@app.route('/videos/<filename>')
def serve_video(filename):
    """Serve generated videos."""
    return send_from_directory(str(VIDEO_DIR), filename)

@app.route('/documents/<filename>')
def serve_document(filename):
    """Serve generated documents."""
    return send_from_directory(str(DOCUMENT_DIR), filename)

@app.route('/audio/<filename>')
def serve_audio(filename):
    """Serve generated audio files."""
    return send_from_directory(str(AUDIO_DIR), filename)

@app.route('/api/gallery', methods=['GET'])
def get_gallery():
    """Get all generated files."""
    file_type = request.args.get('type')
    limit = int(request.args.get('limit', 50))
    files = get_generated_files(file_type, limit)
    return jsonify(files)

@app.route('/api/gallery/<file_id>', methods=['DELETE'])
def delete_gallery_file(file_id):
    """Delete a generated file."""
    delete_generated_file(file_id)
    return jsonify({'success': True})

@app.route('/static/<path:filename>')
def serve_static(filename):
    """Serve static files."""
    return send_from_directory(str(BASE_DIR / 'static'), filename)

# =============================================================================
# HTML Template
# =============================================================================

HTML_TEMPLATE = '''
<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <meta name="theme-color" content="#8b5cf6">
    <meta name="description" content="Ultimate Chat - AI-powered chat interface">
    <title>Ultimate Chat</title>
    <link rel="manifest" href="/static/manifest.json">
    <link rel="preconnect" href="https://fonts.googleapis.com">
    <link rel="preconnect" href="https://fonts.gstatic.com" crossorigin>
    <link href="https://fonts.googleapis.com/css2?family=JetBrains+Mono:wght@400;500&family=Outfit:wght@300;400;500;600;700&display=swap" rel="stylesheet">
    <link rel="stylesheet" href="https://cdnjs.cloudflare.com/ajax/libs/highlight.js/11.9.0/styles/tokyo-night-dark.min.css">
    <style>
        :root {
            --bg-primary: #0a0a0f;
            --bg-secondary: #12121a;
            --bg-tertiary: #1a1a25;
            --bg-hover: #22222f;
            --text-primary: #f0f0f5;
            --text-secondary: #a0a0b0;
            --text-muted: #606070;
            --accent: #8b5cf6;
            --accent-secondary: #ec4899;
            --accent-tertiary: #3b82f6;
            --border: #2a2a35;
            --border-light: #3a3a45;
            --success: #10b981;
            --warning: #f59e0b;
            --error: #ef4444;
            --gradient-1: linear-gradient(135deg, #8b5cf6 0%, #ec4899 100%);
            --gradient-2: linear-gradient(135deg, #3b82f6 0%, #8b5cf6 100%);
            --gradient-3: linear-gradient(135deg, #06b6d4 0%, #3b82f6 100%);
            --shadow-glow: 0 0 40px rgba(139, 92, 246, 0.15);
            --radius-sm: 8px;
            --radius-md: 12px;
            --radius-lg: 16px;
            --radius-xl: 24px;
        }

        * {
            margin: 0;
            padding: 0;
            box-sizing: border-box;
        }

        body {
            font-family: 'Outfit', -apple-system, BlinkMacSystemFont, sans-serif;
            background: var(--bg-primary);
            color: var(--text-primary);
            height: 100vh;
            overflow: hidden;
        }

        /* Layout */
        .app-container {
            display: flex;
            height: 100vh;
        }

        /* Sidebar */
        .sidebar {
            width: 280px;
            background: var(--bg-secondary);
            border-right: 1px solid var(--border);
            display: flex;
            flex-direction: column;
            transition: transform 0.3s ease;
        }

        .sidebar-header {
            padding: 20px;
            border-bottom: 1px solid var(--border);
            display: flex;
            align-items: center;
            gap: 12px;
        }

        .logo {
            width: 40px;
            height: 40px;
            background: var(--gradient-1);
            border-radius: var(--radius-md);
            display: flex;
            align-items: center;
            justify-content: center;
            font-size: 20px;
        }

        .logo-text {
            font-weight: 600;
            font-size: 18px;
            background: var(--gradient-1);
            -webkit-background-clip: text;
            -webkit-text-fill-color: transparent;
            background-clip: text;
        }

        .new-chat-btn {
            margin: 16px;
            padding: 12px 20px;
            background: var(--gradient-1);
            border: none;
            border-radius: var(--radius-md);
            color: white;
            font-family: inherit;
            font-size: 14px;
            font-weight: 500;
            cursor: pointer;
            display: flex;
            align-items: center;
            justify-content: center;
            gap: 8px;
            transition: transform 0.2s, box-shadow 0.2s;
        }

        .new-chat-btn:hover {
            transform: translateY(-2px);
            box-shadow: var(--shadow-glow);
        }

        /* Sidebar Tabs */
        .sidebar-tabs {
            display: flex;
            padding: 8px;
            gap: 4px;
            border-bottom: 1px solid var(--border);
        }
        
        .sidebar-tab {
            flex: 1;
            padding: 10px 12px;
            background: transparent;
            border: none;
            border-radius: var(--radius-sm);
            color: var(--text-muted);
            font-size: 13px;
            cursor: pointer;
            transition: all 0.2s;
        }
        .sidebar-tab:hover {
            background: var(--bg-tertiary);
            color: var(--text-secondary);
        }
        .sidebar-tab.active {
            background: var(--bg-tertiary);
            color: var(--text-primary);
        }
        
        .sidebar-panel {
            flex: 1;
            overflow-y: auto;
            display: none;
        }
        .sidebar-panel.active {
            display: flex;
            flex-direction: column;
        }
        
        /* Gallery */
        .gallery-filters {
            display: flex;
            padding: 8px;
            gap: 4px;
            flex-wrap: wrap;
        }
        
        .filter-btn {
            padding: 6px 12px;
            background: var(--bg-tertiary);
            border: 1px solid transparent;
            border-radius: var(--radius-sm);
            color: var(--text-muted);
            font-size: 12px;
            cursor: pointer;
            transition: all 0.2s;
        }
        .filter-btn:hover {
            border-color: var(--border);
        }
        .filter-btn.active {
            background: var(--accent);
            color: white;
        }
        
        .gallery-grid {
            display: grid;
            grid-template-columns: repeat(2, 1fr);
            gap: 8px;
            padding: 8px;
            overflow-y: auto;
        }
        
        .gallery-item {
            aspect-ratio: 1;
            border-radius: var(--radius-sm);
            overflow: hidden;
            cursor: pointer;
            position: relative;
            background: var(--bg-tertiary);
            border: 1px solid var(--border);
            transition: all 0.2s;
        }
        .gallery-item:hover {
            border-color: var(--accent);
            transform: scale(1.02);
        }
        
        .gallery-item img {
            width: 100%;
            height: 100%;
            object-fit: cover;
        }
        
        .gallery-item-icon {
            width: 100%;
            height: 100%;
            display: flex;
            align-items: center;
            justify-content: center;
            font-size: 32px;
            background: var(--bg-secondary);
        }
        
        .gallery-item-overlay {
            position: absolute;
            inset: 0;
            background: linear-gradient(to top, rgba(0,0,0,0.8), transparent);
            opacity: 0;
            transition: opacity 0.2s;
            display: flex;
            align-items: flex-end;
            padding: 8px;
        }
        .gallery-item:hover .gallery-item-overlay {
            opacity: 1;
        }
        
        .gallery-item-name {
            font-size: 11px;
            color: white;
            white-space: nowrap;
            overflow: hidden;
            text-overflow: ellipsis;
        }
        
        .gallery-item-delete {
            position: absolute;
            top: 4px;
            right: 4px;
            width: 24px;
            height: 24px;
            background: rgba(239, 68, 68, 0.9);
            border: none;
            border-radius: 50%;
            color: white;
            font-size: 12px;
            cursor: pointer;
            opacity: 0;
            transition: opacity 0.2s;
        }
        .gallery-item:hover .gallery-item-delete {
            opacity: 1;
        }
        
        .gallery-empty {
            grid-column: 1 / -1;
            padding: 40px 20px;
            text-align: center;
            color: var(--text-muted);
        }
        .gallery-empty-icon {
            font-size: 48px;
            margin-bottom: 12px;
        }

        .conversations-list {
            flex: 1;
            overflow-y: auto;
            padding: 8px;
        }

        .conversation-item {
            padding: 12px 16px;
            border-radius: var(--radius-md);
            cursor: pointer;
            display: flex;
            align-items: center;
            gap: 12px;
            margin-bottom: 4px;
            transition: background 0.2s;
            position: relative;
        }

        .conversation-item:hover {
            background: var(--bg-hover);
        }

        .conversation-item.active {
            background: var(--bg-tertiary);
        }

        .conversation-icon {
            width: 32px;
            height: 32px;
            background: var(--bg-tertiary);
            border-radius: var(--radius-sm);
            display: flex;
            align-items: center;
            justify-content: center;
            font-size: 14px;
        }

        .conversation-info {
            flex: 1;
            min-width: 0;
        }

        .conversation-title {
            font-size: 14px;
            font-weight: 500;
            white-space: nowrap;
            overflow: hidden;
            text-overflow: ellipsis;
        }

        .conversation-date {
            font-size: 12px;
            color: var(--text-muted);
        }

        .conversation-delete {
            opacity: 0;
            background: none;
            border: none;
            color: var(--text-muted);
            cursor: pointer;
            padding: 4px;
            font-size: 16px;
            transition: opacity 0.2s, color 0.2s;
        }

        .conversation-item:hover .conversation-delete {
            opacity: 1;
        }

        .conversation-delete:hover {
            color: var(--error);
        }

        /* Main Content */
        .main-content {
            flex: 1;
            display: flex;
            flex-direction: column;
            min-width: 0;
        }

        /* Header */
        .header {
            padding: 16px 24px;
            border-bottom: 1px solid var(--border);
            display: flex;
            align-items: center;
            justify-content: space-between;
            background: var(--bg-secondary);
        }

        .header-title {
            font-size: 16px;
            font-weight: 500;
        }

        .model-selector {
            padding: 8px 14px;
            background: var(--bg-tertiary);
            border: 1px solid var(--border);
            border-radius: var(--radius-md);
            font-size: 13px;
            color: var(--text-primary);
            display: flex;
            align-items: center;
            gap: 8px;
            cursor: pointer;
            position: relative;
            transition: all 0.2s;
        }
        .model-selector:hover { border-color: var(--accent); }
        .model-provider-icon { font-size: 16px; }
        .dropdown-arrow { font-size: 10px; color: var(--text-muted); transition: transform 0.2s; }
        .model-selector.open .dropdown-arrow { transform: rotate(180deg); }
        
        .model-dropdown {
            position: absolute;
            top: calc(100% + 8px);
            right: 0;
            min-width: 280px;
            max-height: 400px;
            overflow-y: auto;
            background: var(--bg-secondary);
            border: 1px solid var(--border);
            border-radius: var(--radius-md);
            box-shadow: 0 10px 40px rgba(0,0,0,0.3);
            display: none;
            z-index: 100;
        }
        .model-selector.open .model-dropdown { display: block; }
        
        .provider-section { padding: 8px 0; border-bottom: 1px solid var(--border); }
        .provider-section:last-child { border-bottom: none; }
        .provider-header { padding: 8px 16px; font-size: 11px; font-weight: 600; color: var(--text-muted); text-transform: uppercase; display: flex; align-items: center; gap: 8px; }
        .provider-header .not-configured { color: var(--warning); font-size: 10px; font-weight: normal; text-transform: none; }
        
        .model-option {
            padding: 10px 16px;
            display: flex;
            align-items: center;
            gap: 10px;
            cursor: pointer;
            transition: background 0.15s;
        }
        .model-option:hover { background: var(--bg-hover); }
        .model-option.active { background: var(--bg-tertiary); }
        .model-option .model-name { flex: 1; font-size: 13px; }
        .model-option .check { color: var(--accent); display: none; }
        .model-option.active .check { display: block; }
        
        /* Generation Toggles */
        .gen-toggles {
            display: flex;
            gap: 6px;
        }
        
        .gen-toggle {
            padding: 6px 12px;
            background: var(--bg-tertiary);
            border: 1px solid var(--border);
            border-radius: var(--radius-sm);
            font-size: 13px;
            cursor: pointer;
            display: flex;
            align-items: center;
            gap: 6px;
            transition: all 0.2s;
            color: var(--text-secondary);
        }
        .gen-toggle:hover {
            border-color: var(--accent);
        }
        .gen-toggle.active {
            background: linear-gradient(135deg, rgba(139, 92, 246, 0.2), rgba(236, 72, 153, 0.2));
            border-color: var(--accent);
            color: var(--text-primary);
        }
        .gen-toggle.disabled {
            opacity: 0.5;
            background: var(--bg-secondary);
        }
        
        .toggle-status {
            font-size: 10px;
            font-weight: 600;
            padding: 2px 6px;
            border-radius: 4px;
            background: var(--success);
            color: white;
        }
        .gen-toggle.disabled .toggle-status {
            background: var(--text-muted);
        }
        
        .settings-btn {
            width: 40px;
            height: 40px;
            background: var(--bg-tertiary);
            border: 1px solid var(--border);
            border-radius: var(--radius-md);
            font-size: 18px;
            cursor: pointer;
            transition: all 0.2s;
        }
        .settings-btn:hover { border-color: var(--accent); transform: rotate(45deg); }
        
        /* Model Status Indicators */
        .model-status {
            font-size: 14px;
            margin-right: 6px;
        }
        
        .model-option.cloud .model-name { color: var(--text-primary); }
        .model-option.downloaded .model-name { color: var(--success); }
        .model-option.not-downloaded { 
            opacity: 0.7;
            cursor: pointer;
        }
        .model-option.not-downloaded:hover {
            opacity: 1;
            background: rgba(139, 92, 246, 0.1);
        }
        .model-option.unavailable {
            opacity: 0.5;
            cursor: not-allowed;
        }
        
        .cloud-badge, .local-badge {
            font-size: 10px;
            padding: 2px 6px;
            border-radius: 4px;
            margin-left: 8px;
        }
        .cloud-badge {
            background: linear-gradient(135deg, #3b82f6, #8b5cf6);
            color: white;
        }
        .local-badge {
            background: var(--bg-tertiary);
            color: var(--text-muted);
            border: 1px solid var(--border);
        }
        
        /* Download Modal */
        .download-modal {
            position: fixed;
            top: 0;
            left: 0;
            width: 100%;
            height: 100%;
            background: rgba(0, 0, 0, 0.8);
            display: flex;
            align-items: center;
            justify-content: center;
            z-index: 10000;
        }
        
        .download-content {
            background: var(--bg-secondary);
            border: 1px solid var(--border);
            border-radius: var(--radius-lg);
            padding: 24px;
            width: 400px;
            max-width: 90%;
        }
        
        .download-header {
            display: flex;
            align-items: center;
            gap: 12px;
            font-size: 18px;
            font-weight: 600;
            margin-bottom: 16px;
        }
        
        .download-icon {
            font-size: 24px;
            animation: bounce 1s infinite;
        }
        
        @keyframes bounce {
            0%, 100% { transform: translateY(0); }
            50% { transform: translateY(-5px); }
        }
        
        .download-status {
            color: var(--text-secondary);
            margin-bottom: 12px;
        }
        
        .download-progress-bar {
            height: 8px;
            background: var(--bg-tertiary);
            border-radius: 4px;
            overflow: hidden;
            margin-bottom: 8px;
        }
        
        .download-progress-fill {
            height: 100%;
            background: linear-gradient(90deg, var(--accent), var(--accent-hover));
            width: 0%;
            transition: width 0.3s;
        }
        
        .download-details {
            font-size: 12px;
            color: var(--text-muted);
            margin-bottom: 16px;
        }
        
        .download-cancel {
            width: 100%;
            padding: 10px;
            background: var(--bg-tertiary);
            border: 1px solid var(--border);
            border-radius: var(--radius-sm);
            color: var(--text-primary);
            cursor: pointer;
            transition: all 0.2s;
        }
        .download-cancel:hover {
            background: rgba(239, 68, 68, 0.2);
            border-color: #ef4444;
        }
        
        /* Modal Styles */
        .modal-overlay {
            position: fixed;
            inset: 0;
            background: rgba(0,0,0,0.8);
            display: flex;
            align-items: center;
            justify-content: center;
            z-index: 1000;
            opacity: 0;
            visibility: hidden;
            transition: all 0.3s;
        }
        .modal-overlay.active { opacity: 1; visibility: visible; }
        
        .modal {
            background: var(--bg-secondary);
            border: 1px solid var(--border);
            border-radius: var(--radius-lg);
            width: 90%;
            max-width: 600px;
            max-height: 80vh;
            overflow: hidden;
            display: flex;
            flex-direction: column;
        }
        
        .modal-header {
            padding: 20px 24px;
            border-bottom: 1px solid var(--border);
            display: flex;
            align-items: center;
            justify-content: space-between;
        }
        .modal-header h2 { font-size: 18px; font-weight: 600; }
        .modal-close { background: none; border: none; color: var(--text-muted); font-size: 24px; cursor: pointer; }
        .modal-close:hover { color: var(--text-primary); }
        
        .modal-body { padding: 24px; overflow-y: auto; }
        
        .settings-tabs { display: flex; gap: 8px; margin-bottom: 20px; }
        .tab-btn {
            padding: 10px 20px;
            background: var(--bg-tertiary);
            border: 1px solid var(--border);
            border-radius: var(--radius-md);
            color: var(--text-secondary);
            font-size: 14px;
            cursor: pointer;
            transition: all 0.2s;
        }
        .tab-btn:hover { border-color: var(--accent); }
        .tab-btn.active { background: var(--accent); border-color: var(--accent); color: white; }
        
        .tab-content { display: none; }
        .tab-content.active { display: block; }
        
        .settings-description { color: var(--text-muted); font-size: 14px; margin-bottom: 20px; }
        
        .provider-list { display: flex; flex-direction: column; gap: 12px; }
        
        .provider-card {
            background: var(--bg-tertiary);
            border: 1px solid var(--border);
            border-radius: var(--radius-md);
            padding: 16px;
        }
        .provider-card-header { display: flex; align-items: center; justify-content: space-between; margin-bottom: 12px; }
        .provider-info { display: flex; align-items: center; gap: 10px; }
        .provider-icon { font-size: 24px; }
        .provider-name { font-weight: 500; }
        .provider-status { font-size: 12px; padding: 4px 8px; border-radius: 4px; }
        .provider-status.configured { background: rgba(16,185,129,0.2); color: var(--success); }
        .provider-status.not-configured { background: rgba(245,158,11,0.2); color: var(--warning); }
        
        .settings-input {
            width: 100%;
            padding: 10px 14px;
            background: var(--bg-secondary);
            border: 1px solid var(--border);
            border-radius: var(--radius-sm);
            color: var(--text-primary);
            font-size: 14px;
            margin-bottom: 8px;
        }
        .settings-input:focus { outline: none; border-color: var(--accent); }
        
        .provider-card-actions { display: flex; gap: 8px; }
        
        .btn-primary {
            padding: 10px 20px;
            background: var(--gradient-1);
            border: none;
            border-radius: var(--radius-sm);
            color: white;
            font-size: 14px;
            cursor: pointer;
            transition: transform 0.2s;
        }
        .btn-primary:hover { transform: translateY(-2px); }
        
        .btn-danger {
            padding: 10px 20px;
            background: var(--bg-secondary);
            border: 1px solid var(--error);
            border-radius: var(--radius-sm);
            color: var(--error);
            font-size: 14px;
            cursor: pointer;
        }
        .btn-danger:hover { background: var(--error); color: white; }
        
        .btn-test {
            padding: 10px 20px;
            background: var(--bg-secondary);
            border: 1px solid var(--accent-tertiary);
            border-radius: var(--radius-sm);
            color: var(--accent-tertiary);
            font-size: 14px;
            cursor: pointer;
            transition: all 0.2s;
        }
        .btn-test:hover { background: var(--accent-tertiary); color: white; }
        .btn-test:disabled { opacity: 0.6; cursor: wait; }
        
        .test-result {
            padding: 0;
            margin-bottom: 8px;
            font-size: 13px;
            border-radius: var(--radius-sm);
            overflow: hidden;
            transition: all 0.3s;
        }
        .test-result.success {
            padding: 10px 14px;
            background: rgba(16, 185, 129, 0.15);
            border: 1px solid var(--success);
            color: var(--success);
        }
        .test-result.error {
            padding: 10px 14px;
            background: rgba(239, 68, 68, 0.15);
            border: 1px solid var(--error);
            color: var(--error);
        }
        
        .model-count {
            font-size: 11px;
            color: var(--text-muted);
            padding: 2px 8px;
            background: var(--bg-secondary);
            border-radius: 10px;
        }
        
        .add-model-form { display: flex; gap: 8px; margin-bottom: 16px; flex-wrap: wrap; }
        .add-model-form select, .add-model-form input { flex: 1; min-width: 150px; margin-bottom: 0; }
        
        .custom-models-list { display: flex; flex-direction: column; gap: 8px; }
        .custom-model-item {
            display: flex;
            align-items: center;
            justify-content: space-between;
            padding: 12px 16px;
            background: var(--bg-tertiary);
            border-radius: var(--radius-sm);
        }
        .custom-model-info { display: flex; align-items: center; gap: 12px; }
        .custom-model-delete { background: none; border: none; color: var(--text-muted); cursor: pointer; font-size: 18px; }
        .custom-model-delete:hover { color: var(--error); }

        @keyframes pulse {
            0%, 100% { opacity: 1; }
            50% { opacity: 0.5; }
        }

        /* Chat Area */
        .chat-area {
            flex: 1;
            overflow-y: auto;
            padding: 24px;
            scroll-behavior: smooth;
        }

        .welcome-screen {
            display: flex;
            flex-direction: column;
            align-items: center;
            justify-content: center;
            height: 100%;
            text-align: center;
            padding: 40px;
        }

        .welcome-icon {
            width: 80px;
            height: 80px;
            background: var(--gradient-1);
            border-radius: var(--radius-xl);
            display: flex;
            align-items: center;
            justify-content: center;
            font-size: 40px;
            margin-bottom: 24px;
            box-shadow: var(--shadow-glow);
        }

        .welcome-title {
            font-size: 32px;
            font-weight: 600;
            margin-bottom: 12px;
            background: var(--gradient-1);
            -webkit-background-clip: text;
            -webkit-text-fill-color: transparent;
            background-clip: text;
        }

        .welcome-subtitle {
            font-size: 16px;
            color: var(--text-secondary);
            margin-bottom: 24px;
            max-width: 500px;
        }
        
        .capabilities-grid {
            display: flex;
            flex-wrap: wrap;
            gap: 12px;
            justify-content: center;
            margin-bottom: 32px;
            max-width: 600px;
        }
        
        .capability-item {
            display: flex;
            align-items: center;
            gap: 8px;
            padding: 8px 16px;
            background: var(--bg-secondary);
            border: 1px solid var(--border);
            border-radius: 20px;
            font-size: 13px;
            color: var(--text-secondary);
            transition: all 0.2s;
        }
        
        .capability-item.highlight {
            background: linear-gradient(135deg, rgba(139, 92, 246, 0.2), rgba(236, 72, 153, 0.2));
            border-color: var(--accent);
            color: var(--text-primary);
        }
        
        .capability-icon {
            font-size: 16px;
        }
        
        .welcome-hint {
            font-size: 14px;
            color: var(--text-muted);
            margin-bottom: 16px;
        }

        .quick-prompts {
            display: grid;
            grid-template-columns: repeat(2, 1fr);
            gap: 12px;
            max-width: 600px;
        }

        .quick-prompt {
            padding: 16px 20px;
            background: var(--bg-secondary);
            border: 1px solid var(--border);
            border-radius: var(--radius-md);
            cursor: pointer;
            text-align: left;
            transition: all 0.2s;
        }

        .quick-prompt:hover {
            border-color: var(--accent);
            transform: translateY(-2px);
        }
        
        .quick-prompt.featured {
            background: linear-gradient(135deg, rgba(139, 92, 246, 0.15), rgba(236, 72, 153, 0.15));
            border-color: var(--accent);
        }
        .quick-prompt.featured:hover {
            background: linear-gradient(135deg, rgba(139, 92, 246, 0.25), rgba(236, 72, 153, 0.25));
            box-shadow: 0 4px 20px rgba(139, 92, 246, 0.2);
        }

        .quick-prompt-icon {
            font-size: 20px;
            margin-bottom: 8px;
        }

        .quick-prompt-text {
            font-size: 14px;
            color: var(--text-secondary);
        }
        
        .quick-prompt.featured .quick-prompt-text {
            color: var(--text-primary);
        }

        /* Messages */
        .message {
            display: flex;
            gap: 16px;
            margin-bottom: 24px;
            animation: fadeIn 0.3s ease;
        }

        @keyframes fadeIn {
            from { opacity: 0; transform: translateY(10px); }
            to { opacity: 1; transform: translateY(0); }
        }

        .message-avatar {
            width: 36px;
            height: 36px;
            border-radius: var(--radius-md);
            display: flex;
            align-items: center;
            justify-content: center;
            font-size: 16px;
            flex-shrink: 0;
        }

        .message.user .message-avatar {
            background: var(--gradient-2);
        }

        .message.assistant .message-avatar {
            background: var(--gradient-1);
        }

        .message-content {
            flex: 1;
            min-width: 0;
        }

        .message-header {
            display: flex;
            align-items: center;
            gap: 8px;
            margin-bottom: 8px;
        }

        .message-author {
            font-weight: 500;
            font-size: 14px;
        }

        .message-time {
            font-size: 12px;
            color: var(--text-muted);
        }

        .message-body {
            font-size: 15px;
            line-height: 1.7;
            color: var(--text-secondary);
        }

        .message-body p {
            margin-bottom: 12px;
        }

        .message-body p:last-child {
            margin-bottom: 0;
        }

        .message-body code {
            font-family: 'JetBrains Mono', monospace;
            background: var(--bg-tertiary);
            padding: 2px 6px;
            border-radius: 4px;
            font-size: 13px;
        }

        .message-body pre {
            background: var(--bg-tertiary);
            border-radius: var(--radius-md);
            padding: 16px;
            overflow-x: auto;
            margin: 12px 0;
            border: 1px solid var(--border);
        }

        .message-body pre code {
            background: none;
            padding: 0;
        }

        .message-body ul, .message-body ol {
            margin: 12px 0;
            padding-left: 24px;
        }

        .message-body li {
            margin-bottom: 4px;
        }

        .message-body blockquote {
            border-left: 3px solid var(--accent);
            padding-left: 16px;
            margin: 12px 0;
            color: var(--text-muted);
        }

        .message-body a {
            color: var(--accent);
            text-decoration: none;
        }

        .message-body a:hover {
            text-decoration: underline;
        }

        .message-body h1, .message-body h2, .message-body h3 {
            color: var(--text-primary);
            margin: 16px 0 8px;
        }

        /* Typing Indicator */
        .typing-indicator {
            display: flex;
            gap: 4px;
            padding: 8px 0;
        }

        .typing-dot {
            width: 8px;
            height: 8px;
            background: var(--accent);
            border-radius: 50%;
            animation: typingBounce 1.4s infinite ease-in-out;
        }

        .typing-dot:nth-child(1) { animation-delay: 0s; }
        .typing-dot:nth-child(2) { animation-delay: 0.2s; }
        .typing-dot:nth-child(3) { animation-delay: 0.4s; }

        @keyframes typingBounce {
            0%, 80%, 100% { transform: translateY(0); }
            40% { transform: translateY(-6px); }
        }

        /* Enhanced Progress Card */
        .progress-card {
            background: var(--bg-tertiary);
            border: 1px solid var(--border);
            border-radius: var(--radius-md);
            padding: 20px;
            margin: 12px 0;
            animation: fadeIn 0.3s ease;
        }

        .progress-header {
            display: flex;
            align-items: center;
            justify-content: space-between;
            margin-bottom: 16px;
        }

        .progress-header-left {
            display: flex;
            align-items: center;
            gap: 12px;
        }

        .progress-icon {
            font-size: 24px;
        }

        .progress-title {
            font-weight: 500;
            font-size: 14px;
        }

        .progress-cancel-btn {
            padding: 6px 12px;
            background: transparent;
            border: 1px solid var(--error);
            border-radius: var(--radius-sm);
            color: var(--error);
            font-size: 12px;
            cursor: pointer;
            display: flex;
            align-items: center;
            gap: 6px;
            transition: all 0.2s;
        }
        .progress-cancel-btn:hover {
            background: var(--error);
            color: white;
        }
        .progress-cancel-btn:disabled {
            opacity: 0.5;
            cursor: not-allowed;
        }

        .progress-bar-container {
            margin-bottom: 16px;
        }

        .progress-bar-bg {
            height: 8px;
            background: var(--bg-secondary);
            border-radius: 4px;
            overflow: hidden;
        }

        .progress-bar-fill {
            height: 100%;
            background: var(--gradient-1);
            border-radius: 4px;
            transition: width 0.3s ease;
            position: relative;
        }

        .progress-bar-fill::after {
            content: '';
            position: absolute;
            top: 0;
            left: 0;
            right: 0;
            bottom: 0;
            background: linear-gradient(90deg, transparent, rgba(255,255,255,0.3), transparent);
            animation: shimmer 1.5s infinite;
        }

        @keyframes shimmer {
            0% { transform: translateX(-100%); }
            100% { transform: translateX(100%); }
        }

        .progress-percentage {
            font-size: 12px;
            color: var(--text-muted);
            margin-top: 6px;
            display: flex;
            justify-content: space-between;
        }

        .progress-steps {
            display: flex;
            flex-direction: column;
            gap: 10px;
        }

        .progress-step {
            display: flex;
            align-items: center;
            gap: 12px;
            font-size: 13px;
            color: var(--text-muted);
            transition: all 0.2s;
        }

        .progress-step.active {
            color: var(--text-primary);
        }

        .progress-step.complete {
            color: var(--success);
        }

        .progress-step.cancelled {
            color: var(--warning);
            text-decoration: line-through;
        }

        .step-indicator {
            width: 24px;
            height: 24px;
            border-radius: 50%;
            display: flex;
            align-items: center;
            justify-content: center;
            font-size: 12px;
            background: var(--bg-secondary);
            border: 2px solid var(--border);
            flex-shrink: 0;
            transition: all 0.2s;
        }

        .progress-step.active .step-indicator {
            border-color: var(--accent);
            background: var(--accent);
            color: white;
        }

        .progress-step.active .step-indicator::before {
            content: '';
            position: absolute;
            width: 24px;
            height: 24px;
            border-radius: 50%;
            border: 2px solid var(--accent);
            animation: ping 1s infinite;
        }

        @keyframes ping {
            0% { transform: scale(1); opacity: 1; }
            100% { transform: scale(1.5); opacity: 0; }
        }

        .progress-step.complete .step-indicator {
            border-color: var(--success);
            background: var(--success);
            color: white;
        }

        .step-spinner {
            width: 16px;
            height: 16px;
            border: 2px solid transparent;
            border-top-color: white;
            border-radius: 50%;
            animation: spin 0.8s linear infinite;
        }

        @keyframes spin {
            to { transform: rotate(360deg); }
        }

        .progress-status {
            margin-top: 12px;
            padding: 10px 14px;
            border-radius: var(--radius-sm);
            font-size: 13px;
            display: none;
        }

        .progress-status.error {
            display: block;
            background: rgba(239, 68, 68, 0.15);
            border: 1px solid var(--error);
            color: var(--error);
        }

        .progress-status.cancelled {
            display: block;
            background: rgba(245, 158, 11, 0.15);
            border: 1px solid var(--warning);
            color: var(--warning);
        }

        .progress-status.complete {
            display: block;
            background: rgba(16, 185, 129, 0.15);
            border: 1px solid var(--success);
            color: var(--success);
        }

        .progress-elapsed {
            font-size: 11px;
            color: var(--text-muted);
            margin-top: 8px;
        }

        /* Generated Media */
        .generated-media-container {
            margin: 16px 0;
            background: var(--bg-tertiary);
            border-radius: var(--radius-md);
            overflow: hidden;
            border: 1px solid var(--border);
        }

        .generated-image {
            width: 100%;
            max-width: 512px;
            cursor: pointer;
            transition: transform 0.2s;
        }

        .generated-image:hover {
            transform: scale(1.02);
        }

        .generated-video {
            width: 100%;
            max-width: 512px;
            border-radius: var(--radius-md);
        }

        .media-actions {
            display: flex;
            gap: 8px;
            padding: 12px;
            background: var(--bg-secondary);
        }

        .media-action-btn {
            padding: 8px 16px;
            background: var(--bg-tertiary);
            border: 1px solid var(--border);
            border-radius: var(--radius-sm);
            color: var(--text-secondary);
            font-size: 13px;
            cursor: pointer;
            text-decoration: none;
            transition: all 0.2s;
            font-family: inherit;
        }

        .media-action-btn:hover {
            background: var(--bg-hover);
            color: var(--text-primary);
        }

        /* Document Link */
        .document-link {
            display: inline-flex;
            align-items: center;
            gap: 8px;
            padding: 12px 20px;
            background: var(--bg-tertiary);
            border: 1px solid var(--border);
            border-radius: var(--radius-md);
            color: var(--text-primary);
            text-decoration: none;
            transition: all 0.2s;
        }

        .document-link:hover {
            border-color: var(--accent);
            transform: translateY(-2px);
        }

        .document-icon {
            font-size: 24px;
        }
        
        /* Generated File Card */
        .generated-file-card {
            display: flex;
            align-items: center;
            gap: 16px;
            padding: 16px 20px;
            background: var(--bg-tertiary);
            border: 1px solid var(--border);
            border-radius: var(--radius-md);
            margin: 12px 0;
            transition: all 0.2s;
        }
        .generated-file-card:hover {
            border-color: var(--accent);
            box-shadow: 0 4px 20px rgba(139, 92, 246, 0.1);
        }
        
        .file-icon {
            font-size: 32px;
            width: 48px;
            height: 48px;
            display: flex;
            align-items: center;
            justify-content: center;
            background: var(--bg-secondary);
            border-radius: var(--radius-sm);
        }
        
        .file-info {
            flex: 1;
            display: flex;
            flex-direction: column;
            gap: 4px;
        }
        
        .file-name {
            font-weight: 500;
            color: var(--text-primary);
        }
        
        .file-type {
            font-size: 12px;
            color: var(--text-muted);
        }
        
        .file-download-btn {
            padding: 10px 16px;
            background: var(--gradient-1);
            border: none;
            border-radius: var(--radius-sm);
            color: white;
            font-size: 14px;
            text-decoration: none;
            cursor: pointer;
            transition: transform 0.2s;
        }
        .file-download-btn:hover {
            transform: scale(1.05);
        }
        
        /* Audio Card */
        .generated-audio-card {
            display: flex;
            align-items: center;
            gap: 16px;
            padding: 16px 20px;
            background: var(--bg-tertiary);
            border: 1px solid var(--border);
            border-radius: var(--radius-md);
            margin: 12px 0;
        }
        
        .audio-icon {
            font-size: 28px;
            width: 48px;
            height: 48px;
            display: flex;
            align-items: center;
            justify-content: center;
            background: linear-gradient(135deg, #10b981 0%, #3b82f6 100%);
            border-radius: 50%;
        }
        
        .audio-player-wrapper {
            flex: 1;
        }
        
        .audio-player {
            width: 100%;
            height: 40px;
            border-radius: 20px;
        }
        
        .audio-player::-webkit-media-controls-panel {
            background: var(--bg-secondary);
        }

        /* Input Area */
        .input-area {
            padding: 20px 24px;
            border-top: 1px solid var(--border);
            background: var(--bg-secondary);
        }

        .drop-zone {
            border: 2px dashed var(--border);
            border-radius: var(--radius-md);
            padding: 16px;
            text-align: center;
            margin-bottom: 12px;
            transition: all 0.2s;
            cursor: pointer;
            display: none;
        }

        .drop-zone.visible {
            display: block;
        }

        .drop-zone.dragover {
            border-color: var(--accent);
            background: rgba(139, 92, 246, 0.1);
        }

        .drop-zone-text {
            color: var(--text-muted);
            font-size: 14px;
        }

        .attachments-preview {
            display: flex;
            gap: 8px;
            flex-wrap: wrap;
            margin-bottom: 12px;
        }

        .attachment-item {
            display: flex;
            align-items: center;
            gap: 8px;
            padding: 8px 12px;
            background: var(--bg-tertiary);
            border-radius: var(--radius-sm);
            font-size: 13px;
        }

        .attachment-remove {
            background: none;
            border: none;
            color: var(--text-muted);
            cursor: pointer;
            font-size: 16px;
        }

        .attachment-remove:hover {
            color: var(--error);
        }

        .input-wrapper {
            display: flex;
            gap: 12px;
            align-items: flex-end;
        }

        .input-container {
            flex: 1;
            background: var(--bg-tertiary);
            border: 1px solid var(--border);
            border-radius: var(--radius-md);
            padding: 12px 16px;
            display: flex;
            align-items: flex-end;
            gap: 12px;
            transition: border-color 0.2s;
        }

        .input-container:focus-within {
            border-color: var(--accent);
        }

        .input-actions {
            display: flex;
            gap: 8px;
        }

        .input-action-btn {
            background: none;
            border: none;
            color: var(--text-muted);
            cursor: pointer;
            font-size: 20px;
            padding: 4px;
            transition: color 0.2s;
        }

        .input-action-btn:hover {
            color: var(--text-primary);
        }

        #message-input {
            flex: 1;
            background: none;
            border: none;
            color: var(--text-primary);
            font-family: inherit;
            font-size: 15px;
            resize: none;
            max-height: 150px;
            min-height: 24px;
            outline: none;
        }

        #message-input::placeholder {
            color: var(--text-muted);
        }

        .send-btn {
            width: 48px;
            height: 48px;
            background: var(--gradient-1);
            border: none;
            border-radius: var(--radius-md);
            color: white;
            font-size: 20px;
            cursor: pointer;
            display: flex;
            align-items: center;
            justify-content: center;
            transition: transform 0.2s, box-shadow 0.2s;
        }

        .send-btn:hover {
            transform: scale(1.05);
            box-shadow: var(--shadow-glow);
        }

        .send-btn:disabled {
            opacity: 0.5;
            cursor: not-allowed;
            transform: none;
        }

        /* Lightbox */
        .lightbox {
            position: fixed;
            inset: 0;
            background: rgba(0, 0, 0, 0.95);
            display: flex;
            align-items: center;
            justify-content: center;
            z-index: 1000;
            opacity: 0;
            visibility: hidden;
            transition: all 0.3s;
            cursor: pointer;
        }

        .lightbox.active {
            opacity: 1;
            visibility: visible;
        }

        .lightbox img {
            max-width: 90%;
            max-height: 90%;
            border-radius: var(--radius-md);
            box-shadow: 0 20px 60px rgba(0, 0, 0, 0.5);
        }

        .lightbox-close {
            position: absolute;
            top: 20px;
            right: 20px;
            background: var(--bg-secondary);
            border: none;
            color: white;
            width: 40px;
            height: 40px;
            border-radius: 50%;
            font-size: 24px;
            cursor: pointer;
            display: flex;
            align-items: center;
            justify-content: center;
        }

        /* Toast */
        .toast-container {
            position: fixed;
            bottom: 24px;
            right: 24px;
            display: flex;
            flex-direction: column;
            gap: 12px;
            z-index: 1001;
        }

        .toast {
            padding: 16px 24px;
            background: var(--bg-secondary);
            border: 1px solid var(--border);
            border-radius: var(--radius-md);
            color: var(--text-primary);
            font-size: 14px;
            display: flex;
            align-items: center;
            gap: 12px;
            animation: slideIn 0.3s ease;
            box-shadow: 0 10px 40px rgba(0, 0, 0, 0.3);
        }

        .toast.success { border-left: 3px solid var(--success); }
        .toast.error { border-left: 3px solid var(--error); }
        .toast.warning { border-left: 3px solid var(--warning); }

        @keyframes slideIn {
            from { transform: translateX(100%); opacity: 0; }
            to { transform: translateX(0); opacity: 1; }
        }

        /* Mobile Menu */
        .mobile-menu-btn {
            display: none;
            background: none;
            border: none;
            color: var(--text-primary);
            font-size: 24px;
            cursor: pointer;
        }

        /* Scrollbar */
        ::-webkit-scrollbar {
            width: 8px;
            height: 8px;
        }

        ::-webkit-scrollbar-track {
            background: var(--bg-primary);
        }

        ::-webkit-scrollbar-thumb {
            background: var(--border);
            border-radius: 4px;
        }

        ::-webkit-scrollbar-thumb:hover {
            background: var(--border-light);
        }

        /* Responsive */
        @media (max-width: 768px) {
            .sidebar {
                position: fixed;
                left: 0;
                top: 0;
                height: 100%;
                z-index: 100;
                transform: translateX(-100%);
            }

            .sidebar.open {
                transform: translateX(0);
            }

            .mobile-menu-btn {
                display: block;
            }

            .quick-prompts {
                grid-template-columns: 1fr;
            }
        }
    </style>
</head>
<body>
    <div class="app-container">
        <!-- Sidebar -->
        <aside class="sidebar" id="sidebar">
            <div class="sidebar-header">
                <div class="logo">ü§ñ</div>
                <span class="logo-text">Ultimate Chat</span>
            </div>
            
            <button class="new-chat-btn" onclick="createNewChat()">
                <span>‚ú®</span> New Chat
            </button>
            
            <!-- Sidebar Tabs -->
            <div class="sidebar-tabs">
                <button class="sidebar-tab active" onclick="switchSidebarTab('chats')" data-tab="chats">
                    üí¨ Chats
                </button>
                <button class="sidebar-tab" onclick="switchSidebarTab('gallery')" data-tab="gallery">
                    üñºÔ∏è Gallery
                </button>
            </div>
            
            <!-- Chats Panel -->
            <div class="sidebar-panel active" id="panel-chats">
                <div class="conversations-list" id="conversations-list">
                    <!-- Conversations loaded dynamically -->
                </div>
            </div>
            
            <!-- Gallery Panel -->
            <div class="sidebar-panel" id="panel-gallery">
                <div class="gallery-filters">
                    <button class="filter-btn active" onclick="filterGallery(null)">All</button>
                    <button class="filter-btn" onclick="filterGallery('image')">üñºÔ∏è</button>
                    <button class="filter-btn" onclick="filterGallery('audio')">üîä</button>
                    <button class="filter-btn" onclick="filterGallery('pdf')">üìÑ</button>
                    <button class="filter-btn" onclick="filterGallery('chart')">üìä</button>
                </div>
                <div class="gallery-grid" id="gallery-grid">
                    <!-- Gallery items loaded dynamically -->
                </div>
            </div>
        </aside>

        <!-- Main Content -->
        <main class="main-content">
            <!-- Header -->
            <header class="header">
                <div style="display: flex; align-items: center; gap: 12px;">
                    <button class="mobile-menu-btn" onclick="toggleSidebar()">‚ò∞</button>
                    <h1 class="header-title" id="header-title">New Conversation</h1>
                </div>
                <div style="display: flex; align-items: center; gap: 12px;">
                    <!-- Generation Toggles -->
                    <div class="gen-toggles">
                        <button class="gen-toggle" id="toggle-image" onclick="toggleGeneration('image')" title="Image Generation">
                            üé® <span class="toggle-status" id="image-status">ON</span>
                        </button>
                        <button class="gen-toggle" id="toggle-video" onclick="toggleGeneration('video')" title="Video Generation">
                            üé¨ <span class="toggle-status" id="video-status">ON</span>
                        </button>
                    </div>
                    <div class="model-selector" onclick="toggleModelDropdown(event)">
                        <span class="model-provider-icon" id="model-provider-icon">ü¶ô</span>
                        <span id="current-model-display">llama3.2:latest</span>
                        <span class="dropdown-arrow">‚ñº</span>
                        <div class="model-dropdown" id="model-dropdown">
                            <!-- Populated by JS -->
                        </div>
                    </div>
                    <button class="settings-btn" onclick="openSettings()" title="Settings">‚öôÔ∏è</button>
                </div>
            </header>

            <!-- Chat Area -->
            <div class="chat-area" id="chat-area">
                <!-- Welcome Screen -->
                <div class="welcome-screen" id="welcome-screen">
                    <div class="welcome-icon">üöÄ</div>
                    <h2 class="welcome-title">Welcome to Ultimate Chat</h2>
                    <p class="welcome-subtitle">
                        Your all-in-one AI assistant. Chat with multiple AI models, generate images, audio, documents, and more.
                    </p>
                    
                    <!-- Capabilities Grid -->
                    <div class="capabilities-grid">
                        <div class="capability-item">
                            <span class="capability-icon">üí¨</span>
                            <span>Multi-Model Chat</span>
                        </div>
                        <div class="capability-item highlight">
                            <span class="capability-icon">üé®</span>
                            <span>SDXL Images</span>
                        </div>
                        <div class="capability-item highlight">
                            <span class="capability-icon">üé¨</span>
                            <span>AI Video</span>
                        </div>
                        <div class="capability-item">
                            <span class="capability-icon">üîä</span>
                            <span>Text to Speech</span>
                        </div>
                        <div class="capability-item">
                            <span class="capability-icon">üìä</span>
                            <span>Charts & Graphs</span>
                        </div>
                        <div class="capability-item">
                            <span class="capability-icon">üì±</span>
                            <span>QR Codes</span>
                        </div>
                        <div class="capability-item">
                            <span class="capability-icon">üìÑ</span>
                            <span>PDF/Word/Excel</span>
                        </div>
                    </div>
                    
                    <p class="welcome-hint">Try these commands to get started:</p>
                    
                    <div class="quick-prompts">
                        <div class="quick-prompt featured" onclick="sendQuickPrompt('Generate image of a cyberpunk city at night with neon lights and flying cars')">
                            <div class="quick-prompt-icon">üé®</div>
                            <div class="quick-prompt-text">Generate AI image (SDXL)</div>
                        </div>
                        <div class="quick-prompt featured" onclick="sendQuickPrompt('Generate video of a serene mountain landscape with flowing clouds')">
                            <div class="quick-prompt-icon">üé¨</div>
                            <div class="quick-prompt-text">Generate AI video</div>
                        </div>
                        <div class="quick-prompt" onclick="sendQuickPrompt('Generate audio: Welcome to Ultimate Chat, your all-in-one AI assistant!')">
                            <div class="quick-prompt-icon">üîä</div>
                            <div class="quick-prompt-text">Generate audio</div>
                        </div>
                        <div class="quick-prompt" onclick="sendQuickPrompt('Create a bar chart with data: Sales:150, Marketing:80, Development:200, Support:60')">
                            <div class="quick-prompt-icon">üìä</div>
                            <div class="quick-prompt-text">Create a chart</div>
                        </div>
                        <div class="quick-prompt" onclick="sendQuickPrompt('Generate QR code for https://github.com')">
                            <div class="quick-prompt-icon">üì±</div>
                            <div class="quick-prompt-text">Generate QR code</div>
                        </div>
                        <div class="quick-prompt" onclick="sendQuickPrompt('Explain how transformers work in AI')">
                            <div class="quick-prompt-icon">üß†</div>
                            <div class="quick-prompt-text">Ask anything</div>
                        </div>
                    </div>
                </div>

                <!-- Messages Container -->
                <div id="messages-container"></div>
            </div>

            <!-- Input Area -->
            <div class="input-area">
                <div class="drop-zone" id="drop-zone" onclick="document.getElementById('file-input').click()">
                    <span class="drop-zone-text">üìé Drop files here or click to upload</span>
                </div>
                
                <div class="attachments-preview" id="attachments-preview"></div>
                
                <div class="input-wrapper">
                    <div class="input-container">
                        <div class="input-actions">
                            <button class="input-action-btn" onclick="toggleDropZone()" title="Attach file">üìé</button>
                        </div>
                        <textarea id="message-input" placeholder="Type your message..." rows="1" onkeydown="handleKeyDown(event)"></textarea>
                    </div>
                    <button class="send-btn" id="send-btn" onclick="sendMessage()">‚û§</button>
                </div>
                
                <input type="file" id="file-input" multiple hidden onchange="handleFileSelect(event)">
            </div>
        </main>
    </div>

    <!-- Lightbox -->
    <div class="lightbox" id="lightbox" onclick="closeLightbox()">
        <button class="lightbox-close">√ó</button>
        <img id="lightbox-image" src="" alt="">
    </div>

    <!-- Toast Container -->
    <div class="toast-container" id="toast-container"></div>

    <!-- Settings Modal -->
    <div class="modal-overlay" id="settings-modal" onclick="if(event.target===this)closeSettings()">
        <div class="modal">
            <div class="modal-header">
                <h2>‚öôÔ∏è Settings</h2>
                <button class="modal-close" onclick="closeSettings()">√ó</button>
            </div>
            <div class="modal-body">
                <div class="settings-tabs">
                    <button class="tab-btn active" onclick="switchTab('api-keys')">API Keys</button>
                    <button class="tab-btn" onclick="switchTab('custom-models')">Custom Models</button>
                </div>
                
                <div class="tab-content active" id="tab-api-keys">
                    <p class="settings-description">Configure your API keys to enable different AI providers.</p>
                    
                    <div class="provider-list" id="provider-list">
                        <!-- Populated by JS -->
                    </div>
                </div>
                
                <div class="tab-content" id="tab-custom-models">
                    <p class="settings-description">Add custom models to any provider.</p>
                    
                    <div class="add-model-form">
                        <select id="custom-model-provider" class="settings-input">
                            <option value="openai">OpenAI</option>
                            <option value="anthropic">Anthropic</option>
                            <option value="google">Google</option>
                            <option value="xai">xAI</option>
                            <option value="ollama">Ollama</option>
                        </select>
                        <input type="text" id="custom-model-id" class="settings-input" placeholder="Model ID (e.g., gpt-4-turbo-2024)">
                        <button class="btn-primary" onclick="addCustomModel()">Add Model</button>
                    </div>
                    
                    <div class="custom-models-list" id="custom-models-list">
                        <!-- Populated by JS -->
                    </div>
                </div>
            </div>
        </div>
    </div>

    <!-- External Scripts -->
    <script src="https://cdn.jsdelivr.net/npm/marked/marked.min.js"></script>
    <script src="https://cdnjs.cloudflare.com/ajax/libs/highlight.js/11.9.0/highlight.min.js"></script>
    
    <script>
        // =============================================================================
        // State
        // =============================================================================
        
        let currentConversationId = null;
        let attachments = [];
        let isStreaming = false;
        let providers = {};
        let currentProvider = 'openai';
        let currentModel = 'gpt-4o-mini';
        let imageGenEnabled = true;
        let videoGenEnabled = true;

        // =============================================================================
        // Initialization
        // =============================================================================
        
        document.addEventListener('DOMContentLoaded', async () => {
            await loadProviders();
            await loadSettings();
            loadConversations();
            setupDragDrop();
            autoResizeTextarea();
            
            document.addEventListener('click', (e) => {
                if (!e.target.closest('.model-selector')) {
                    document.querySelector('.model-selector')?.classList.remove('open');
                }
            });
            
            if ('serviceWorker' in navigator) {
                navigator.serviceWorker.register('/static/service-worker.js').catch(() => {});
            }
        });

        // =============================================================================
        // Provider & Model Management
        // =============================================================================
        
        let ollamaModels = { available: false, downloaded: [] };
        let isDownloading = false;
        
        async function loadProviders() {
            try {
                const response = await fetch('/api/providers');
                providers = await response.json();
                await loadOllamaStatus();
                renderModelDropdown();
            } catch (e) { console.error('Failed to load providers:', e); }
        }
        
        async function loadOllamaStatus() {
            try {
                const response = await fetch('/api/ollama/models');
                ollamaModels = await response.json();
            } catch (e) {
                ollamaModels = { available: false, downloaded: [] };
            }
        }
        
        async function loadSettings() {
            try {
                const response = await fetch('/api/settings');
                const settings = await response.json();
                currentProvider = settings.current_provider || 'openai';
                currentModel = settings.current_model || 'gpt-4o-mini';
                updateModelDisplay();
            } catch (e) { console.error('Failed to load settings:', e); }
        }
        
        function getModelStatus(provider, model) {
            if (provider !== 'ollama') {
                return { icon: '‚òÅÔ∏è', status: 'cloud', tooltip: 'Cloud model' };
            }
            if (!ollamaModels.available) {
                return { icon: '‚ö†Ô∏è', status: 'unavailable', tooltip: 'Ollama not running' };
            }
            const modelBase = model.split(':')[0];
            const isDownloaded = ollamaModels.downloaded.some(m => m.includes(modelBase));
            if (isDownloaded) {
                return { icon: '‚úÖ', status: 'downloaded', tooltip: 'Downloaded locally' };
            }
            return { icon: 'üì•', status: 'not-downloaded', tooltip: 'Click to download' };
        }
        
        function renderModelDropdown() {
            const dropdown = document.getElementById('model-dropdown');
            let html = '';
            
            // Show cloud providers first (recommended)
            const sortedProviders = Object.entries(providers).sort((a, b) => {
                if (a[0] === 'ollama') return 1;
                if (b[0] === 'ollama') return -1;
                return 0;
            });
            
            for (const [pid, pdata] of sortedProviders) {
                const isCloud = pid !== 'ollama';
                html += `<div class="provider-section ${isCloud ? 'cloud-provider' : 'local-provider'}">
                    <div class="provider-header">
                        <span>${pdata.icon} ${pdata.name}</span>
                        ${isCloud ? '<span class="cloud-badge">‚òÅÔ∏è Cloud</span>' : '<span class="local-badge">üíª Local</span>'}
                        ${!pdata.configured && pid !== 'ollama' ? '<span class="not-configured">‚ö†Ô∏è No API key</span>' : ''}
                    </div>`;
                
                for (const model of pdata.models) {
                    const isActive = pid === currentProvider && model === currentModel;
                    const status = getModelStatus(pid, model);
                    const needsDownload = status.status === 'not-downloaded';
                    
                    html += `<div class="model-option ${isActive ? 'active' : ''} ${status.status}" 
                                  onclick="${needsDownload ? `promptDownload('${model}')` : `selectModel('${pid}', '${model}')`}"
                                  data-provider="${pid}" data-model="${model}"
                                  title="${status.tooltip}">
                        <span class="model-status">${status.icon}</span>
                        <span class="model-name">${model}</span>
                        ${isActive ? '<span class="check">‚úì</span>' : ''}
                    </div>`;
                }
                html += '</div>';
            }
            dropdown.innerHTML = html;
        }
        
        function toggleModelDropdown(e) {
            e.stopPropagation();
            document.querySelector('.model-selector').classList.toggle('open');
        }
        
        async function selectModel(provider, model) {
            currentProvider = provider;
            currentModel = model;
            updateModelDisplay();
            document.querySelector('.model-selector').classList.remove('open');
            
            await fetch('/api/settings/model', {
                method: 'POST',
                headers: {'Content-Type': 'application/json'},
                body: JSON.stringify({provider, model})
            });
            
            renderModelDropdown();
            showToast(`Switched to ${model}`, 'success');
        }
        
        function updateModelDisplay() {
            const icon = providers[currentProvider]?.icon || 'ü§ñ';
            document.getElementById('model-provider-icon').textContent = icon;
            document.getElementById('current-model-display').textContent = currentModel.length > 25 ? currentModel.slice(0,22) + '...' : currentModel;
        }
        
        // =============================================================================
        // Settings Modal
        // =============================================================================
        
        function openSettings() {
            document.getElementById('settings-modal').classList.add('active');
            renderProviderList();
            renderCustomModelsList();
        }
        
        function closeSettings() {
            document.getElementById('settings-modal').classList.remove('active');
        }
        
        function switchTab(tabId) {
            document.querySelectorAll('.tab-btn').forEach(btn => btn.classList.remove('active'));
            document.querySelectorAll('.tab-content').forEach(tab => tab.classList.remove('active'));
            document.querySelector(`[onclick="switchTab('${tabId}')"]`).classList.add('active');
            document.getElementById('tab-' + tabId).classList.add('active');
        }
        
        async function renderProviderList() {
            const settings = await (await fetch('/api/settings')).json();
            const container = document.getElementById('provider-list');
            
            let html = '';
            for (const [pid, pdata] of Object.entries(providers)) {
                const apiKeyInfo = settings.api_keys[pid];
                const configured = apiKeyInfo?.configured;
                
                html += `<div class="provider-card" id="provider-card-${pid}">
                    <div class="provider-card-header">
                        <div class="provider-info">
                            <span class="provider-icon">${pdata.icon}</span>
                            <span class="provider-name">${pdata.name}</span>
                            <span class="model-count">${pdata.models.length} models</span>
                        </div>
                        <span class="provider-status ${configured ? 'configured' : 'not-configured'}" id="status-${pid}">
                            ${configured ? '‚úì Configured' : 'Not configured'}
                        </span>
                    </div>
                    <div class="test-result" id="test-result-${pid}"></div>
                    <input type="password" class="settings-input" id="api-key-${pid}" 
                           placeholder="${configured ? 'Enter new API key to update...' : 'Enter API key...'}"
                           ${configured ? 'data-has-key="true"' : ''}>
                    ${pid === 'ollama' ? `<input type="text" class="settings-input" id="base-url-${pid}" 
                           placeholder="Base URL (default: http://localhost:11434)" 
                           value="${apiKeyInfo?.base_url || ''}">` : ''}
                    <div class="provider-card-actions">
                        <button class="btn-primary" onclick="saveApiKey('${pid}')">Save</button>
                        <button class="btn-test" onclick="testConnection('${pid}')" id="test-btn-${pid}">
                            üîå Test Connection
                        </button>
                        ${configured ? `<button class="btn-danger" onclick="deleteApiKey('${pid}')">Remove</button>` : ''}
                    </div>
                </div>`;
            }
            container.innerHTML = html;
        }
        
        async function testConnection(provider) {
            const btn = document.getElementById(`test-btn-${provider}`);
            const resultDiv = document.getElementById(`test-result-${provider}`);
            
            btn.disabled = true;
            btn.innerHTML = '‚è≥ Testing...';
            resultDiv.innerHTML = '';
            resultDiv.className = 'test-result';
            
            try {
                const response = await fetch(`/api/test/${provider}`, {method: 'POST'});
                const data = await response.json();
                
                if (data.success) {
                    resultDiv.innerHTML = `‚úÖ ${data.message}`;
                    resultDiv.className = 'test-result success';
                    showToast(`${providers[provider].name}: Connection successful!`, 'success');
                } else {
                    resultDiv.innerHTML = `‚ùå ${data.error}`;
                    resultDiv.className = 'test-result error';
                    showToast(`${providers[provider].name}: ${data.error}`, 'error');
                }
            } catch (e) {
                resultDiv.innerHTML = `‚ùå Connection failed: ${e.message}`;
                resultDiv.className = 'test-result error';
                showToast('Test failed', 'error');
            } finally {
                btn.disabled = false;
                btn.innerHTML = 'üîå Test Connection';
            }
        }
        
        async function saveApiKey(provider) {
            const apiKey = document.getElementById(`api-key-${provider}`).value;
            const baseUrlEl = document.getElementById(`base-url-${provider}`);
            const baseUrl = baseUrlEl?.value || null;
            
            if (!apiKey && !baseUrl) {
                showToast('Please enter an API key', 'warning');
                return;
            }
            
            await fetch('/api/settings/apikey', {
                method: 'POST',
                headers: {'Content-Type': 'application/json'},
                body: JSON.stringify({provider, api_key: apiKey || 'local', base_url: baseUrl})
            });
            
            await loadProviders();
            renderProviderList();
            renderModelDropdown();
            showToast(`${providers[provider].name} API key saved`, 'success');
        }
        
        async function deleteApiKey(provider) {
            if (!confirm(`Remove ${providers[provider].name} API key?`)) return;
            
            await fetch(`/api/settings/apikey/${provider}`, {method: 'DELETE'});
            await loadProviders();
            renderProviderList();
            renderModelDropdown();
            showToast('API key removed', 'success');
        }
        
        async function renderCustomModelsList() {
            const settings = await (await fetch('/api/settings')).json();
            const container = document.getElementById('custom-models-list');
            
            if (!settings.custom_models.length) {
                container.innerHTML = '<p style="color: var(--text-muted); text-align: center;">No custom models added yet.</p>';
                return;
            }
            
            container.innerHTML = settings.custom_models.map(m => `
                <div class="custom-model-item">
                    <div class="custom-model-info">
                        <span>${providers[m.provider]?.icon || 'ü§ñ'}</span>
                        <span>${m.model_id}</span>
                        <span style="color: var(--text-muted); font-size: 12px;">(${m.provider})</span>
                    </div>
                    <button class="custom-model-delete" onclick="deleteCustomModel('${m.id}')">üóëÔ∏è</button>
                </div>
            `).join('');
        }
        
        async function addCustomModel() {
            const provider = document.getElementById('custom-model-provider').value;
            const modelId = document.getElementById('custom-model-id').value.trim();
            
            if (!modelId) {
                showToast('Please enter a model ID', 'warning');
                return;
            }
            
            await fetch('/api/models/custom', {
                method: 'POST',
                headers: {'Content-Type': 'application/json'},
                body: JSON.stringify({provider, model_id: modelId})
            });
            
            document.getElementById('custom-model-id').value = '';
            await loadProviders();
            renderModelDropdown();
            renderCustomModelsList();
            showToast('Custom model added', 'success');
        }
        
        async function deleteCustomModel(id) {
            await fetch(`/api/models/custom/${id}`, {method: 'DELETE'});
            await loadProviders();
            renderModelDropdown();
            renderCustomModelsList();
            showToast('Model removed', 'success');
        }

        // =============================================================================
        // Conversations
        // =============================================================================
        
        async function loadConversations() {
            try {
                const response = await fetch('/api/conversations');
                const conversations = await response.json();
                renderConversations(conversations);
            } catch (error) {
                console.error('Failed to load conversations:', error);
            }
        }

        function renderConversations(conversations) {
            const container = document.getElementById('conversations-list');
            container.innerHTML = conversations.map(conv => `
                <div class="conversation-item ${conv.id === currentConversationId ? 'active' : ''}" 
                     onclick="loadConversation('${conv.id}')"
                     data-id="${conv.id}">
                    <div class="conversation-icon">üí¨</div>
                    <div class="conversation-info">
                        <div class="conversation-title">${escapeHtml(conv.title)}</div>
                        <div class="conversation-date">${formatDate(conv.updated_at)}</div>
                    </div>
                    <button class="conversation-delete" onclick="event.stopPropagation(); deleteConversation('${conv.id}')">üóëÔ∏è</button>
                </div>
            `).join('');
        }

        async function createNewChat() {
            try {
                const response = await fetch('/api/conversation/create', { method: 'POST' });
                const data = await response.json();
                currentConversationId = data.id;
                
                document.getElementById('welcome-screen').style.display = 'flex';
                document.getElementById('messages-container').innerHTML = '';
                document.getElementById('header-title').textContent = 'New Conversation';
                
                loadConversations();
                closeSidebar();
            } catch (error) {
                showToast('Failed to create conversation', 'error');
            }
        }

        async function loadConversation(id) {
            try {
                currentConversationId = id;
                const response = await fetch(`/api/conversation/${id}`);
                const messages = await response.json();
                
                document.getElementById('welcome-screen').style.display = 'none';
                renderMessages(messages);
                
                // Update active state
                document.querySelectorAll('.conversation-item').forEach(item => {
                    item.classList.toggle('active', item.dataset.id === id);
                });
                
                // Update title
                const conv = document.querySelector(`.conversation-item[data-id="${id}"] .conversation-title`);
                if (conv) {
                    document.getElementById('header-title').textContent = conv.textContent;
                }
                
                closeSidebar();
            } catch (error) {
                showToast('Failed to load conversation', 'error');
            }
        }

        async function deleteConversation(id) {
            if (!confirm('Delete this conversation?')) return;
            
            try {
                await fetch(`/api/conversation/${id}`, { method: 'DELETE' });
                
                if (currentConversationId === id) {
                    currentConversationId = null;
                    document.getElementById('welcome-screen').style.display = 'flex';
                    document.getElementById('messages-container').innerHTML = '';
                    document.getElementById('header-title').textContent = 'New Conversation';
                }
                
                loadConversations();
                showToast('Conversation deleted', 'success');
            } catch (error) {
                showToast('Failed to delete conversation', 'error');
            }
        }

        // =============================================================================
        // Messages
        // =============================================================================
        
        function renderMessages(messages) {
            const container = document.getElementById('messages-container');
            container.innerHTML = messages.map(msg => createMessageHTML(msg.role, msg.content, msg.attachments)).join('');
            scrollToBottom();
            hljs.highlightAll();
        }

        function createMessageHTML(role, content, msgAttachments = []) {
            const isUser = role === 'user';
            const avatar = isUser ? 'üë§' : 'ü§ñ';
            const author = isUser ? 'You' : 'Assistant';
            const time = formatTime(new Date());
            
            let attachmentsHTML = '';
            if (msgAttachments && msgAttachments.length > 0) {
                attachmentsHTML = msgAttachments.map(att => {
                    if (att.type && att.type.startsWith('image/')) {
                        return `<img src="${att.url}" alt="${att.filename}" style="max-width: 200px; border-radius: 8px; margin-top: 8px;">`;
                    }
                    return `<div class="attachment-item">${att.filename}</div>`;
                }).join('');
            }
            
            const renderedContent = isUser ? escapeHtml(content) : marked.parse(content);
            
            return `
                <div class="message ${role}">
                    <div class="message-avatar">${avatar}</div>
                    <div class="message-content">
                        <div class="message-header">
                            <span class="message-author">${author}</span>
                            <span class="message-time">${time}</span>
                        </div>
                        <div class="message-body">${renderedContent}</div>
                        ${attachmentsHTML}
                    </div>
                </div>
            `;
        }

        function addMessage(role, content, msgAttachments = []) {
            const container = document.getElementById('messages-container');
            container.insertAdjacentHTML('beforeend', createMessageHTML(role, content, msgAttachments));
            scrollToBottom();
            if (role === 'assistant') {
                hljs.highlightAll();
            }
        }

        function updateLastAssistantMessage(content) {
            const messages = document.querySelectorAll('.message.assistant');
            if (messages.length > 0) {
                const lastMessage = messages[messages.length - 1];
                const body = lastMessage.querySelector('.message-body');
                body.innerHTML = marked.parse(content);
                scrollToBottom();
            }
        }

        // =============================================================================
        // Send Message
        // =============================================================================
        
        async function sendMessage() {
            const input = document.getElementById('message-input');
            const message = input.value.trim();
            
            if (!message && attachments.length === 0) return;
            if (isStreaming) return;
            
            // Create conversation if needed
            if (!currentConversationId) {
                const response = await fetch('/api/conversation/create', { method: 'POST' });
                const data = await response.json();
                currentConversationId = data.id;
            }
            
            // Hide welcome screen
            document.getElementById('welcome-screen').style.display = 'none';
            
            // Add user message
            addMessage('user', message, attachments);
            
            // Clear input
            input.value = '';
            input.style.height = 'auto';
            const currentAttachments = [...attachments];
            attachments = [];
            document.getElementById('attachments-preview').innerHTML = '';
            
            // Start streaming
            isStreaming = true;
            document.getElementById('send-btn').disabled = true;
            
            // Add assistant message placeholder with typing indicator
            const container = document.getElementById('messages-container');
            container.insertAdjacentHTML('beforeend', `
                <div class="message assistant" id="streaming-message">
                    <div class="message-avatar">ü§ñ</div>
                    <div class="message-content">
                        <div class="message-header">
                            <span class="message-author">Assistant</span>
                            <span class="message-time">${formatTime(new Date())}</span>
                        </div>
                        <div class="message-body">
                            <div class="typing-indicator">
                                <div class="typing-dot"></div>
                                <div class="typing-dot"></div>
                                <div class="typing-dot"></div>
                            </div>
                        </div>
                    </div>
                </div>
            `);
            scrollToBottom();
            
            try {
                const response = await fetch('/api/chat/stream', {
                    method: 'POST',
                    headers: { 'Content-Type': 'application/json' },
                    body: JSON.stringify({
                        message,
                        conversation_id: currentConversationId,
                        attachments: currentAttachments,
                        provider: currentProvider,
                        model: currentModel,
                        image_gen_enabled: imageGenEnabled,
                        video_gen_enabled: videoGenEnabled
                    })
                });
                
                const reader = response.body.getReader();
                const decoder = new TextDecoder();
                let fullContent = '';
                
                while (true) {
                    const { done, value } = await reader.read();
                    if (done) break;
                    
                    const chunk = decoder.decode(value);
                    const lines = chunk.split('\\n');
                    
                    for (const line of lines) {
                        if (line.startsWith('data: ')) {
                            try {
                                const data = JSON.parse(line.slice(6));
                                
                                // Handle task lifecycle events
                                if (data.task_start) {
                                    showProgressCard(data.task_start);
                                }
                                
                                if (data.task_progress) {
                                    updateProgressCard(data.task_progress);
                                }
                                
                                if (data.task_complete) {
                                    completeProgressCard(data.task_complete);
                                }
                                
                                if (data.content) {
                                    fullContent += data.content;
                                    const streamingMsg = document.getElementById('streaming-message');
                                    if (streamingMsg) {
                                        const body = streamingMsg.querySelector('.message-body');
                                        body.innerHTML = marked.parse(fullContent);
                                    }
                                    scrollToBottom();
                                }
                                
                                if (data.image) {
                                    fullContent += `\\n\\n<div class="generated-media-container">
                                        <img src="${data.image}" class="generated-image" onclick="openLightbox('${data.image}')">
                                        <div class="media-actions">
                                            <button onclick="openLightbox('${data.image}')" class="media-action-btn">üîç View Full Size</button>
                                            <a href="${data.image}" target="_blank" class="media-action-btn">üîó Open in Tab</a>
                                            <a href="${data.image}" download class="media-action-btn">üíæ Download</a>
                                        </div>
                                    </div>`;
                                    const streamingMsg = document.getElementById('streaming-message');
                                    if (streamingMsg) {
                                        const body = streamingMsg.querySelector('.message-body');
                                        body.innerHTML = marked.parse(fullContent);
                                    }
                                }
                                
                                if (data.video) {
                                    fullContent += `\\n\\n<div class="generated-media-container">
                                        <video src="${data.video}" class="generated-video" controls></video>
                                        <div class="media-actions">
                                            <a href="${data.video}" target="_blank" class="media-action-btn">üîó Open in Tab</a>
                                            <a href="${data.video}" download class="media-action-btn">üíæ Download</a>
                                        </div>
                                    </div>`;
                                    const streamingMsg = document.getElementById('streaming-message');
                                    if (streamingMsg) {
                                        const body = streamingMsg.querySelector('.message-body');
                                        body.innerHTML = marked.parse(fullContent);
                                    }
                                }
                                
                                if (data.document) {
                                    const ext = (data.document_name || '').split('.').pop().toLowerCase();
                                    const icons = {pdf: 'üìï', docx: 'üìò', pptx: 'üìô', xlsx: 'üìó'};
                                    const icon = icons[ext] || 'üìÑ';
                                    fullContent += `\\n\\n<div class="generated-file-card">
                                        <div class="file-icon">${icon}</div>
                                        <div class="file-info">
                                            <span class="file-name">${data.document_name || 'Document'}</span>
                                            <span class="file-type">${ext.toUpperCase()} Document</span>
                                        </div>
                                        <a href="${data.document}" download class="file-download-btn">‚¨áÔ∏è Download</a>
                                    </div>`;
                                    const streamingMsg = document.getElementById('streaming-message');
                                    if (streamingMsg) {
                                        const body = streamingMsg.querySelector('.message-body');
                                        body.innerHTML = marked.parse(fullContent);
                                    }
                                }
                                
                                if (data.audio) {
                                    fullContent += `\\n\\n<div class="generated-audio-card">
                                        <div class="audio-icon">üîä</div>
                                        <div class="audio-player-wrapper">
                                            <audio controls class="audio-player" src="${data.audio}">
                                                Your browser does not support audio.
                                            </audio>
                                        </div>
                                        <a href="${data.audio}" download="${data.audio_name || 'audio.mp3'}" class="file-download-btn">‚¨áÔ∏è</a>
                                    </div>`;
                                    const streamingMsg = document.getElementById('streaming-message');
                                    if (streamingMsg) {
                                        const body = streamingMsg.querySelector('.message-body');
                                        body.innerHTML = marked.parse(fullContent);
                                    }
                                }
                                
                                if (data.done) {
                                    const streamingMsg = document.getElementById('streaming-message');
                                    if (streamingMsg) {
                                        streamingMsg.removeAttribute('id');
                                    }
                                    hljs.highlightAll();
                                    loadConversations();
                                }
                            } catch (e) {
                                console.error('Parse error:', e);
                            }
                        }
                    }
                }
            } catch (error) {
                console.error('Stream error:', error);
                showToast('Failed to get response', 'error');
            } finally {
                isStreaming = false;
                document.getElementById('send-btn').disabled = false;
            }
        }

        function sendQuickPrompt(prompt) {
            document.getElementById('message-input').value = prompt;
            sendMessage();
        }

        // =============================================================================
        // File Upload
        // =============================================================================
        
        function setupDragDrop() {
            const dropZone = document.getElementById('drop-zone');
            const chatArea = document.getElementById('chat-area');
            
            ['dragenter', 'dragover', 'dragleave', 'drop'].forEach(eventName => {
                chatArea.addEventListener(eventName, preventDefaults, false);
            });
            
            chatArea.addEventListener('dragenter', () => dropZone.classList.add('visible', 'dragover'));
            chatArea.addEventListener('dragleave', (e) => {
                if (!chatArea.contains(e.relatedTarget)) {
                    dropZone.classList.remove('dragover');
                }
            });
            chatArea.addEventListener('drop', handleDrop);
        }

        function preventDefaults(e) {
            e.preventDefault();
            e.stopPropagation();
        }

        function handleDrop(e) {
            const dropZone = document.getElementById('drop-zone');
            dropZone.classList.remove('visible', 'dragover');
            
            const files = e.dataTransfer.files;
            handleFiles(files);
        }

        function handleFileSelect(e) {
            handleFiles(e.target.files);
        }

        async function handleFiles(files) {
            for (const file of files) {
                try {
                    const formData = new FormData();
                    formData.append('file', file);
                    
                    const response = await fetch('/api/upload', {
                        method: 'POST',
                        body: formData
                    });
                    
                    const data = await response.json();
                    
                    attachments.push({
                        url: data.url,
                        filename: data.filename,
                        type: file.type
                    });
                    
                    updateAttachmentsPreview();
                    showToast(`Uploaded ${file.name}`, 'success');
                } catch (error) {
                    showToast(`Failed to upload ${file.name}`, 'error');
                }
            }
        }

        function updateAttachmentsPreview() {
            const container = document.getElementById('attachments-preview');
            container.innerHTML = attachments.map((att, i) => `
                <div class="attachment-item">
                    ${att.type.startsWith('image/') ? `<img src="${att.url}" style="width: 32px; height: 32px; object-fit: cover; border-radius: 4px;">` : 'üìé'}
                    <span>${att.filename}</span>
                    <button class="attachment-remove" onclick="removeAttachment(${i})">√ó</button>
                </div>
            `).join('');
        }

        function removeAttachment(index) {
            attachments.splice(index, 1);
            updateAttachmentsPreview();
        }

        function toggleDropZone() {
            const dropZone = document.getElementById('drop-zone');
            dropZone.classList.toggle('visible');
        }

        // =============================================================================
        // Lightbox
        // =============================================================================
        
        function openLightbox(src) {
            const lightbox = document.getElementById('lightbox');
            const image = document.getElementById('lightbox-image');
            image.src = src;
            lightbox.classList.add('active');
        }

        function closeLightbox() {
            document.getElementById('lightbox').classList.remove('active');
        }

        // =============================================================================
        // Toast
        // =============================================================================
        
        function showToast(message, type = 'info') {
            const container = document.getElementById('toast-container');
            const toast = document.createElement('div');
            toast.className = `toast ${type}`;
            toast.innerHTML = `
                ${type === 'success' ? '‚úì' : type === 'error' ? '‚úï' : '‚Ñπ'}
                <span>${message}</span>
            `;
            container.appendChild(toast);
            
            setTimeout(() => {
                toast.style.animation = 'slideIn 0.3s ease reverse';
                setTimeout(() => toast.remove(), 300);
            }, 3000);
        }

        // =============================================================================
        // Enhanced Progress Cards with Cancellation
        // =============================================================================
        
        let activeTaskId = null;
        let taskStartTime = null;
        let taskElapsedInterval = null;

        function showProgressCard(taskData) {
            const container = document.getElementById('messages-container');
            activeTaskId = taskData.id;
            taskStartTime = Date.now();
            
            const icon = taskData.type === 'image' ? 'üé®' : taskData.type === 'video' ? 'üé¨' : 'üìÑ';
            
            container.insertAdjacentHTML('beforeend', `
                <div class="progress-card" id="progress-card-${taskData.id}">
                    <div class="progress-header">
                        <div class="progress-header-left">
                            <span class="progress-icon">${icon}</span>
                            <span class="progress-title">${taskData.title}</span>
                        </div>
                        <button class="progress-cancel-btn" id="cancel-btn-${taskData.id}" onclick="cancelTask('${taskData.id}')">
                            ‚úï Cancel
                        </button>
                    </div>
                    <div class="progress-bar-container">
                        <div class="progress-bar-bg">
                            <div class="progress-bar-fill" id="progress-bar-${taskData.id}" style="width: 0%"></div>
                        </div>
                        <div class="progress-percentage">
                            <span id="progress-text-${taskData.id}">Starting...</span>
                            <span id="progress-percent-${taskData.id}">0%</span>
                        </div>
                    </div>
                    <div class="progress-steps" id="progress-steps-${taskData.id}">
                        ${taskData.steps.map((step, i) => `
                            <div class="progress-step" id="step-${taskData.id}-${i}" data-step="${i}">
                                <div class="step-indicator">${i + 1}</div>
                                <span>${step}</span>
                            </div>
                        `).join('')}
                    </div>
                    <div class="progress-status" id="progress-status-${taskData.id}"></div>
                    <div class="progress-elapsed" id="progress-elapsed-${taskData.id}">Elapsed: 0s</div>
                </div>
            `);
            
            // Start elapsed time counter
            taskElapsedInterval = setInterval(() => {
                const elapsed = Math.floor((Date.now() - taskStartTime) / 1000);
                const elapsedEl = document.getElementById(`progress-elapsed-${taskData.id}`);
                if (elapsedEl) {
                    const mins = Math.floor(elapsed / 60);
                    const secs = elapsed % 60;
                    elapsedEl.textContent = mins > 0 ? `Elapsed: ${mins}m ${secs}s` : `Elapsed: ${secs}s`;
                }
            }, 1000);
            
            scrollToBottom();
        }

        function updateProgressCard(taskProgress) {
            const taskId = activeTaskId;
            if (!taskId) return;
            
            const progressBar = document.getElementById(`progress-bar-${taskId}`);
            const progressText = document.getElementById(`progress-text-${taskId}`);
            const progressPercent = document.getElementById(`progress-percent-${taskId}`);
            const statusEl = document.getElementById(`progress-status-${taskId}`);
            const cancelBtn = document.getElementById(`cancel-btn-${taskId}`);
            
            if (progressBar) progressBar.style.width = `${taskProgress.progress || 0}%`;
            if (progressPercent) progressPercent.textContent = `${taskProgress.progress || 0}%`;
            if (progressText) progressText.textContent = taskProgress.step_name || 'Processing...';
            
            // Update steps
            const currentStep = taskProgress.step || 0;
            const stepsContainer = document.getElementById(`progress-steps-${taskId}`);
            if (stepsContainer) {
                stepsContainer.querySelectorAll('.progress-step').forEach((stepEl, i) => {
                    stepEl.classList.remove('active', 'complete', 'cancelled');
                    const indicator = stepEl.querySelector('.step-indicator');
                    
                    if (i < currentStep) {
                        stepEl.classList.add('complete');
                        indicator.innerHTML = '‚úì';
                    } else if (i === currentStep - 1 || (i === 0 && currentStep === 0)) {
                        if (taskProgress.status === 'running' || taskProgress.status === 'starting') {
                            stepEl.classList.add('active');
                            indicator.innerHTML = '<div class="step-spinner"></div>';
                        }
                    } else {
                        indicator.textContent = i + 1;
                    }
                });
            }
            
            // Handle status
            if (taskProgress.status === 'completed') {
                if (statusEl) {
                    statusEl.className = 'progress-status complete';
                    statusEl.textContent = '‚úÖ Completed successfully!';
                }
                if (cancelBtn) cancelBtn.style.display = 'none';
                clearInterval(taskElapsedInterval);
            } else if (taskProgress.status === 'cancelled') {
                if (statusEl) {
                    statusEl.className = 'progress-status cancelled';
                    statusEl.textContent = '‚ö†Ô∏è Operation was cancelled';
                }
                if (cancelBtn) cancelBtn.style.display = 'none';
                clearInterval(taskElapsedInterval);
                // Mark remaining steps as cancelled
                if (stepsContainer) {
                    stepsContainer.querySelectorAll('.progress-step:not(.complete)').forEach(s => {
                        s.classList.add('cancelled');
                    });
                }
            } else if (taskProgress.status === 'error') {
                if (statusEl) {
                    statusEl.className = 'progress-status error';
                    statusEl.textContent = `‚ùå Error: ${taskProgress.error || 'Unknown error'}`;
                }
                if (cancelBtn) cancelBtn.style.display = 'none';
                clearInterval(taskElapsedInterval);
            } else if (taskProgress.status === 'cancelling') {
                if (cancelBtn) {
                    cancelBtn.disabled = true;
                    cancelBtn.textContent = '‚è≥ Cancelling...';
                }
            }
            
            scrollToBottom();
        }

        async function cancelTask(taskId) {
            const cancelBtn = document.getElementById(`cancel-btn-${taskId}`);
            if (cancelBtn) {
                cancelBtn.disabled = true;
                cancelBtn.textContent = '‚è≥ Cancelling...';
            }
            
            try {
                await fetch(`/api/task/${taskId}/cancel`, { method: 'POST' });
                showToast('Cancellation requested...', 'warning');
            } catch (e) {
                console.error('Failed to cancel task:', e);
                showToast('Failed to cancel', 'error');
            }
        }

        function completeProgressCard(taskId) {
            activeTaskId = null;
            clearInterval(taskElapsedInterval);
            
            // Auto-remove card after 3 seconds if successful
            setTimeout(() => {
                const card = document.getElementById(`progress-card-${taskId}`);
                if (card) {
                    card.style.opacity = '0';
                    card.style.transform = 'translateY(-10px)';
                    card.style.transition = 'all 0.3s ease';
                    setTimeout(() => card.remove(), 300);
                }
            }, 3000);
        }

        // Legacy functions for compatibility
        function showProgress(type, title) {
            showProgressCard({ id: 'legacy', type, title, steps: ['Processing...'] });
        }

        function updateProgress(stepIndex) {
            updateProgressCard({ step: stepIndex, progress: stepIndex * 25, status: 'running' });
        }

        function removeProgress() {
            completeProgressCard('legacy');
        }

        // =============================================================================
        // Utilities
        // =============================================================================
        
        function scrollToBottom() {
            const chatArea = document.getElementById('chat-area');
            chatArea.scrollTop = chatArea.scrollHeight;
        }

        function autoResizeTextarea() {
            const textarea = document.getElementById('message-input');
            textarea.addEventListener('input', () => {
                textarea.style.height = 'auto';
                textarea.style.height = Math.min(textarea.scrollHeight, 150) + 'px';
            });
        }

        function handleKeyDown(e) {
            if (e.key === 'Enter' && !e.shiftKey) {
                e.preventDefault();
                sendMessage();
            }
        }

        function escapeHtml(text) {
            const div = document.createElement('div');
            div.textContent = text;
            return div.innerHTML;
        }

        function formatDate(dateStr) {
            const date = new Date(dateStr);
            const now = new Date();
            const diff = now - date;
            
            if (diff < 60000) return 'Just now';
            if (diff < 3600000) return Math.floor(diff / 60000) + 'm ago';
            if (diff < 86400000) return Math.floor(diff / 3600000) + 'h ago';
            return date.toLocaleDateString();
        }

        function formatTime(date) {
            return date.toLocaleTimeString([], { hour: '2-digit', minute: '2-digit' });
        }

        function toggleSidebar() {
            document.getElementById('sidebar').classList.toggle('open');
        }

        function closeSidebar() {
            document.getElementById('sidebar').classList.remove('open');
        }
        
        // =============================================================================
        // Sidebar Tabs & Gallery
        // =============================================================================
        
        let currentGalleryFilter = null;
        
        function switchSidebarTab(tabName) {
            document.querySelectorAll('.sidebar-tab').forEach(t => t.classList.remove('active'));
            document.querySelectorAll('.sidebar-panel').forEach(p => p.classList.remove('active'));
            
            document.querySelector(`.sidebar-tab[data-tab="${tabName}"]`).classList.add('active');
            document.getElementById(`panel-${tabName}`).classList.add('active');
            
            if (tabName === 'gallery') {
                loadGallery();
            }
        }
        
        async function loadGallery(fileType = null) {
            currentGalleryFilter = fileType;
            const grid = document.getElementById('gallery-grid');
            
            try {
                const url = fileType ? `/api/gallery?type=${fileType}` : '/api/gallery';
                const response = await fetch(url);
                const files = await response.json();
                
                if (files.length === 0) {
                    grid.innerHTML = `
                        <div class="gallery-empty">
                            <div class="gallery-empty-icon">üìÅ</div>
                            <div>No files yet</div>
                            <div style="font-size: 12px; margin-top: 8px;">
                                Try generating images, audio, or documents!
                            </div>
                        </div>
                    `;
                    return;
                }
                
                grid.innerHTML = files.map(file => {
                    const isImage = ['image', 'qr', 'chart'].includes(file.file_type);
                    const icons = {
                        audio: 'üîä', pdf: 'üìï', docx: 'üìò', pptx: 'üìô', xlsx: 'üìó', 
                        image: 'üñºÔ∏è', qr: 'üì±', chart: 'üìä'
                    };
                    
                    return `
                        <div class="gallery-item" onclick="openGalleryItem('${file.file_path}', '${file.file_type}')">
                            ${isImage 
                                ? `<img src="${file.file_path}" alt="${file.file_name}" loading="lazy">`
                                : `<div class="gallery-item-icon">${icons[file.file_type] || 'üìÑ'}</div>`
                            }
                            <div class="gallery-item-overlay">
                                <span class="gallery-item-name">${file.file_name}</span>
                            </div>
                            <button class="gallery-item-delete" onclick="event.stopPropagation(); deleteGalleryItem('${file.id}')">√ó</button>
                        </div>
                    `;
                }).join('');
                
            } catch (e) {
                console.error('Failed to load gallery:', e);
                grid.innerHTML = '<div class="gallery-empty">Failed to load gallery</div>';
            }
        }
        
        function filterGallery(fileType) {
            document.querySelectorAll('.filter-btn').forEach(b => b.classList.remove('active'));
            event.target.classList.add('active');
            loadGallery(fileType);
        }
        
        function openGalleryItem(path, type) {
            if (['image', 'qr', 'chart'].includes(type)) {
                openLightbox(path);
            } else if (type === 'audio') {
                // Play audio in a modal or just open
                window.open(path, '_blank');
            } else {
                // Download document
                const a = document.createElement('a');
                a.href = path;
                a.download = '';
                a.click();
            }
        }
        
        async function deleteGalleryItem(fileId) {
            if (!confirm('Delete this file?')) return;
            
            try {
                await fetch(`/api/gallery/${fileId}`, { method: 'DELETE' });
                loadGallery(currentGalleryFilter);
                showToast('File deleted', 'success');
            } catch (e) {
                showToast('Failed to delete file', 'error');
            }
        }
        
        // =============================================================================
        // Generation Toggles
        // =============================================================================
        
        function initGenerationToggles() {
            // Load from localStorage
            imageGenEnabled = localStorage.getItem('imageGenEnabled') !== 'false';
            videoGenEnabled = localStorage.getItem('videoGenEnabled') !== 'false';
            updateToggleUI();
        }
        
        function toggleGeneration(type) {
            if (type === 'image') {
                imageGenEnabled = !imageGenEnabled;
                localStorage.setItem('imageGenEnabled', imageGenEnabled);
            } else if (type === 'video') {
                videoGenEnabled = !videoGenEnabled;
                localStorage.setItem('videoGenEnabled', videoGenEnabled);
            }
            updateToggleUI();
            showToast(`${type.charAt(0).toUpperCase() + type.slice(1)} generation ${type === 'image' ? (imageGenEnabled ? 'enabled' : 'disabled') : (videoGenEnabled ? 'enabled' : 'disabled')}`, 'success');
        }
        
        function updateToggleUI() {
            const imageToggle = document.getElementById('toggle-image');
            const videoToggle = document.getElementById('toggle-video');
            const imageStatus = document.getElementById('image-status');
            const videoStatus = document.getElementById('video-status');
            
            if (imageToggle) {
                imageToggle.classList.toggle('active', imageGenEnabled);
                imageToggle.classList.toggle('disabled', !imageGenEnabled);
                imageStatus.textContent = imageGenEnabled ? 'ON' : 'OFF';
            }
            
            if (videoToggle) {
                videoToggle.classList.toggle('active', videoGenEnabled);
                videoToggle.classList.toggle('disabled', !videoGenEnabled);
                videoStatus.textContent = videoGenEnabled ? 'ON' : 'OFF';
            }
        }
        
        // Initialize toggles on page load
        document.addEventListener('DOMContentLoaded', initGenerationToggles);
        
        // =============================================================================
        // Ollama Model Download
        // =============================================================================
        
        function promptDownload(model) {
            if (isDownloading) {
                showToast('A download is already in progress', 'warning');
                return;
            }
            
            const confirmed = confirm(`Download "${model}" from Ollama?\n\nThis may take several minutes depending on model size and your connection speed.\n\nClick OK to start downloading.`);
            
            if (confirmed) {
                startModelDownload(model);
            }
        }
        
        async function startModelDownload(model) {
            isDownloading = true;
            
            // Show download progress modal
            showDownloadModal(model);
            
            try {
                const eventSource = new EventSource(`/api/ollama/pull/stream/${encodeURIComponent(model)}`);
                
                eventSource.onmessage = (event) => {
                    const data = JSON.parse(event.data);
                    updateDownloadProgress(data);
                    
                    if (data.status === 'success' || data.error) {
                        eventSource.close();
                        isDownloading = false;
                        
                        if (data.status === 'success') {
                            showToast(`Model "${model}" downloaded successfully!`, 'success');
                            setTimeout(() => {
                                hideDownloadModal();
                                loadOllamaStatus().then(() => {
                                    renderModelDropdown();
                                    selectModel('ollama', model);
                                });
                            }, 1000);
                        } else {
                            showToast(`Download failed: ${data.error}`, 'error');
                            hideDownloadModal();
                        }
                    }
                };
                
                eventSource.onerror = () => {
                    eventSource.close();
                    isDownloading = false;
                    showToast('Download connection lost', 'error');
                    hideDownloadModal();
                };
            } catch (e) {
                isDownloading = false;
                showToast(`Download failed: ${e.message}`, 'error');
                hideDownloadModal();
            }
        }
        
        function showDownloadModal(model) {
            // Remove existing modal if any
            const existing = document.getElementById('download-modal');
            if (existing) existing.remove();
            
            const modal = document.createElement('div');
            modal.id = 'download-modal';
            modal.className = 'download-modal';
            modal.innerHTML = `
                <div class="download-content">
                    <div class="download-header">
                        <span class="download-icon">üì•</span>
                        <span>Downloading ${model}</span>
                    </div>
                    <div class="download-status" id="download-status">Connecting to Ollama...</div>
                    <div class="download-progress-bar">
                        <div class="download-progress-fill" id="download-progress-fill"></div>
                    </div>
                    <div class="download-details" id="download-details"></div>
                    <button class="download-cancel" onclick="cancelDownload()">Cancel</button>
                </div>
            `;
            document.body.appendChild(modal);
        }
        
        function hideDownloadModal() {
            const modal = document.getElementById('download-modal');
            if (modal) modal.remove();
        }
        
        function updateDownloadProgress(data) {
            const statusEl = document.getElementById('download-status');
            const progressEl = document.getElementById('download-progress-fill');
            const detailsEl = document.getElementById('download-details');
            
            if (!statusEl) return;
            
            if (data.status) {
                statusEl.textContent = data.status;
            }
            
            if (data.completed && data.total) {
                const percent = Math.round((data.completed / data.total) * 100);
                progressEl.style.width = percent + '%';
                const completedMB = (data.completed / 1024 / 1024).toFixed(1);
                const totalMB = (data.total / 1024 / 1024).toFixed(1);
                detailsEl.textContent = `${completedMB} MB / ${totalMB} MB (${percent}%)`;
            } else if (data.status === 'success') {
                progressEl.style.width = '100%';
                statusEl.textContent = '‚úÖ Download complete!';
            }
        }
        
        function cancelDownload() {
            // Note: Can't actually cancel the Ollama download, but we can close the modal
            hideDownloadModal();
            isDownloading = false;
            showToast('Download cancelled (may continue in background)', 'warning');
        }
    </script>
</body>
</html>
'''

# =============================================================================
# Main Entry Point
# =============================================================================

if __name__ == '__main__':
    print(f"""
‚ïî‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïó
‚ïë                    ü§ñ Ultimate Chat                           ‚ïë
‚ïë          Multi-Provider AI Chat Interface                     ‚ïë
‚ï†‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ï£
‚ïë  Server: http://localhost:{APP_PORT}                             ‚ïë
‚ïë  Providers: OpenAI, Anthropic, Google, xAI, Ollama           ‚ïë
‚ïë  Image Generation: {'Enabled' if ENABLE_IMAGE_GEN else 'Disabled':<39}‚ïë
‚ïë  Environment: {'Production' if IS_PRODUCTION else 'Development':<43}‚ïë
‚ïö‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïê‚ïù
    """)
    app.run(host='0.0.0.0', port=APP_PORT, debug=not IS_PRODUCTION, threaded=True)

